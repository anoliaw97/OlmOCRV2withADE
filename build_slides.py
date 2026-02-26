"""
Build new slides for [24] 26 Feb 2026_ver01.pptx  -> ver02
Appends 8 new slides after the existing 13.

Theme (confirmed from theme2.xml + existing slide XML):
  - Background:   WHITE (FFFFFF) - not dark navy
  - Text color:   002060 for headings (existing style), BLACK (000000) for body
  - Accent 1:     4472C4 (blue)
  - Accent 2:     ED7D31 (orange)
  - Accent 4:     FFC000 (yellow)
  - Accent 5:     5B9BD5 (light blue)
  - Accent 6:     70AD47 (green)
  - Fonts:        Calibri Light (headings), Calibri (body)
  - Slide size:   10693400 x 7556500 EMU (A4 landscape)

Rules:
  - Minimum font size: 12pt everywhere
  - Body text color: BLACK (000000), no grey/washed-out colors
  - Headings: 002060 navy (matching existing slides)
  - Table headers: white text on navy background
"""

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── File paths ─────────────────────────────────────────────────────────────────
PPTX_IN  = r"C:/Users/Mining/Downloads/olmocr-main/OlmOCRV2withADE/[24] 26 Feb 2026_ver01.pptx"
PPTX_OUT = r"C:/Users/Mining/Downloads/olmocr-main/OlmOCRV2withADE/[24] 26 Feb 2026_ver05.pptx"

# ── Results PNG directory ────────────────────────────────────────────────────────
RESULTS_DIR = r"C:/Users/Mining/Downloads/olmocr-main/OlmOCRV2withADE/RESULTS"


# ══════════════════════════════════════════════════════════════════════════════
#  IMAGE HELPER
# ══════════════════════════════════════════════════════════════════════════════

def add_png_shapes(slide, left, top, width, height, png_paths):
    """Add PNG images, scaling to fit."""
    for i, png_path in enumerate(png_paths):
        try:
            img_left = left + i * (width + Emu(20_000))
            slide.shapes.add_picture(png_path, img_left, top, width=width, height=height)
        except Exception as e:
            print(f"  Warning: Could not add {png_path}: {e}")

# ── Colors (matching real theme) ───────────────────────────────────────────────
C_NAVY   = RGBColor(0x00, 0x20, 0x60)   # 002060 – heading / title color (used in existing slides)
C_BLUE   = RGBColor(0x44, 0x72, 0xC4)   # 4472C4 – accent 1
C_ORANGE = RGBColor(0xED, 0x7D, 0x31)   # ED7D31 – accent 2
C_YELLOW = RGBColor(0xFF, 0xC0, 0x00)   # FFC000 – accent 4
C_LBLUE  = RGBColor(0x5B, 0x9B, 0xD5)   # 5B9BD5 – accent 5
C_GREEN  = RGBColor(0x70, 0xAD, 0x47)   # 70AD47 – accent 6
C_BLACK  = RGBColor(0x00, 0x00, 0x00)   # 000000 – body text
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # light grey – table row alternating
C_MGRAY  = RGBColor(0xD6, 0xDC, 0xE4)   # medium grey – borders

# Slide dimensions
SW = 10693400
SH = 7556500

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_blank_slide(prs):
    for layout in prs.slide_master.slide_layouts:
        if layout.name == "Blank":
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_master.slide_layouts[4])


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0.75):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def txb(slide, left, top, width, height, text, size=14, bold=False,
        color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Simple single-run textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(max(size, 12))   # enforce minimum 12pt
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _esc(t):
    return t.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")


def para_xml(text, size=14, bold=False, color="000000", bullet=None,
             indent=0, underline=False, align="l", italic=False):
    """Build a raw <a:p> XML element."""
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sz = int(max(size, 12) * 100)
    b = ' b="1"' if bold else ''
    u = ' u="sng"' if underline else ''
    it = ' i="1"' if italic else ''
    fill = f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'

    mar = f' marL="{indent*457200}" indent="-457200"' if indent else ''
    if bullet == "auto":
        bxml = '<a:buFont typeface="Calibri"/><a:buAutoNum type="arabicPeriod"/>'
    elif bullet == "dot":
        bxml = '<a:buChar char="•"/>'
    elif bullet == "dash":
        bxml = '<a:buChar char="–"/>'
    elif bullet == "check":
        bxml = '<a:buChar char="✓"/>'
    elif bullet == "cross":
        bxml = '<a:buChar char="✗"/>'
    elif bullet == "circle":
        bxml = '<a:buChar char="○"/>'
    else:
        bxml = '<a:buNone/>'

    xml = (
        f'<a:p xmlns:a="{a}">'
        f'<a:pPr algn="{align}"{mar}>{bxml}</a:pPr>'
        f'<a:r><a:rPr lang="en-US" sz="{sz}"{b}{u}{it} dirty="0">'
        f'{fill}</a:rPr><a:t>{_esc(text)}</a:t></a:r>'
        f'</a:p>'
    )
    return etree.fromstring(xml)


def para_empty(size=12):
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sz = int(max(size, 12) * 100)
    return etree.fromstring(
        f'<a:p xmlns:a="{a}"><a:endParaRPr lang="en-US" sz="{sz}" dirty="0"/></a:p>'
    )


def multiline_tb(slide, left, top, width, height, paras, wrap=True):
    """Textbox from list of <a:p> XML elements."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    txBody = tf._txBody
    # set anchor top
    bpr = txBody.find(qn("a:bodyPr"))
    if bpr is not None:
        bpr.set("anchor", "t")
    for old in txBody.findall(qn("a:p")):
        txBody.remove(old)
    for p in paras:
        txBody.append(p)
    return box


def title_bar(slide, title, subtitle=None):
    """Navy top bar matching existing slide style (title bold, white)."""
    bar_h = Emu(1_180_000)
    add_rect(slide, 0, 0, SW, bar_h, fill_rgb=C_NAVY)
    # Title
    txb(slide, Emu(350_000), Emu(180_000), Emu(9_800_000), Emu(600_000),
        title, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, Emu(350_000), Emu(780_000), Emu(9_800_000), Emu(380_000),
            subtitle, size=13, bold=False, color=RGBColor(0xBF,0xD7,0xED),
            align=PP_ALIGN.LEFT)
    return bar_h


def styled_table(slide, left, top, width, height, headers, rows,
                 col_ratios=None, hdr_bg=C_NAVY, hdr_fg=C_WHITE,
                 body_size=12, hdr_size=13):
    """
    Add a table. col_ratios: list summing to 1.0.
    Even rows: LGRAY, odd rows: WHITE.
    All body text: BLACK, minimum 12pt.
    """
    ncols = len(headers)
    nrows = len(rows) + 1
    ts = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = ts.table

    if col_ratios is None:
        col_ratios = [1/ncols] * ncols
    for ci, r in enumerate(col_ratios):
        tbl.columns[ci].width = int(width * r)

    # Header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(max(hdr_size, 12))
        run.font.color.rgb = hdr_fg
        run.font.name = "Calibri"

    # Data rows
    for ri, row in enumerate(rows):
        bg = C_LGRAY if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(max(body_size, 12))
            run.font.bold = (ci == 0)
            run.font.color.rgb = C_BLACK
            run.font.name = "Calibri"

    return tbl


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1: Work Directions (Revised – 4 directions)
# ══════════════════════════════════════════════════════════════════════════════

def slide_work_directions(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Work Directions",
                   "Current & future research roadmap — structured into 4 progressive steps")

    CARD_TOP = bh + Emu(200_000)
    CARD_H   = SH - CARD_TOP - Emu(200_000)
    GAP      = Emu(120_000)
    CARD_W   = (SW - Emu(400_000) - 3*GAP) // 4

    cards = [
        ("1", "Prompt Refinement",
         [
             "Establish stable default prompt using olmOCR v4 (build_no_anchoring_v4_yaml_prompt).",
             "Benchmark: 17-page document ~40 min; 6–8 pages ~30 min.",
             "Greedy decoding (temp = 0.0), max 8000 new tokens.",
         ],
         "Completed", C_GREEN),
        ("2", "Specific Use-Case Prompting",
         [
             "Task-specific prompts for 3 use cases:",
             "  (a) Target SCAL parameters by keyword",
             "  (b) Graph digitization — model limitation identified",
             "  (c) Column-specific extraction — ongoing prompt search",
         ],
         "Ongoing", C_ORANGE),
        ("3", "Agentic Document Extraction",
         [
             "Full ADE pipeline: PDF → olmOCR FP8 → raw text → PostProcessAgent (LLM) → structured dataset → Export.",
             "GUI with prompt optimizer, page timing, table renderer, Excel/CSV/JSON export.",
         ],
         "In Progress", C_BLUE),
        ("4", "Custom Model Fine-Tuning",
         [
             "Build SCAL/rel-perm training corpus from internal PDFs.",
             "Fine-tune olmOCR using grpo_train.py.",
             "Goal: PETRO-olmOCR — domain-adapted model version.",
         ],
         "Planned", C_LBLUE),
    ]

    for i, (num, title, bullets, status, sc) in enumerate(cards):
        lft = Emu(200_000) + i * (CARD_W + GAP)

        # Card border
        add_rect(slide, lft, CARD_TOP, CARD_W, CARD_H,
                 fill_rgb=C_WHITE, line_rgb=C_MGRAY, line_width_pt=1.0)

        # Colored top accent strip
        add_rect(slide, lft, CARD_TOP, CARD_W, Emu(220_000), fill_rgb=sc)

        # Number badge
        add_rect(slide, lft + Emu(60_000), CARD_TOP - Emu(120_000),
                 Emu(450_000), Emu(450_000), fill_rgb=sc)
        txb(slide, lft + Emu(60_000), CARD_TOP - Emu(120_000),
            Emu(450_000), Emu(450_000),
            num, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Title
        txb(slide, lft + Emu(80_000), CARD_TOP + Emu(260_000),
            CARD_W - Emu(160_000), Emu(550_000),
            title, size=16, bold=True, color=C_NAVY, wrap=True)

        # Bullet points
        paras = []
        for b in bullets:
            paras.append(para_xml(b, size=12, color="000000",
                                  bullet="dot" if not b.startswith("  ") else None,
                                  indent=1 if not b.startswith("  ") else 2))
            paras.append(para_empty(6))
        multiline_tb(slide,
                     lft + Emu(80_000), CARD_TOP + Emu(880_000),
                     CARD_W - Emu(160_000), CARD_H - Emu(1_350_000),
                     paras)

        # Status badge at bottom
        add_rect(slide, lft + Emu(80_000), CARD_TOP + CARD_H - Emu(420_000),
                 CARD_W - Emu(160_000), Emu(340_000), fill_rgb=sc)
        txb(slide, lft + Emu(80_000), CARD_TOP + CARD_H - Emu(420_000),
            CARD_W - Emu(160_000), Emu(340_000),
            status, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2: What is ADE?
# ══════════════════════════════════════════════════════════════════════════════

def slide_ade_what_is(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Agentic Document Extraction (ADE)",
                   "LandingAI framework — docs.landing.ai/ade/ade-overview  |  github.com/landing-ai/ade-helper-scripts")

    TOP   = bh + Emu(200_000)
    H     = SH - TOP - Emu(200_000)
    HALF  = SW // 2 - Emu(200_000)
    GAP   = Emu(150_000)

    # ---- Left column: definition + pipeline ----
    add_rect(slide, Emu(150_000), TOP, HALF, H,
             fill_rgb=C_LGRAY, line_rgb=C_MGRAY)

    left_paras = [
        para_xml("What is ADE?", 18, bold=True, color="002060", underline=True),
        para_empty(8),
        para_xml(
            "Agentic Document Extraction (ADE) is a document intelligence pipeline "
            "that converts unstructured documents — PDFs, scans, images — into reliable, "
            "structured, machine-readable data suitable for RAG, search, and ML applications.",
            13, color="000000"),
        para_empty(10),
        para_xml("Core ADE Pipeline (LandingAI):", 14, bold=True, color="002060"),
        para_empty(4),
        para_xml("Parse", 13, bold=True, color="4472C4"),
        para_xml(
            "Converts documents to structured Markdown + JSON chunks. "
            "Each chunk (text, table, figure) gets page number and bounding box coordinates.",
            12, color="000000", bullet="dot", indent=1),
        para_empty(6),
        para_xml("Split  (optional)", 13, bold=True, color="4472C4"),
        para_xml(
            "Separates multi-document files into individual documents by type "
            "(e.g., batch KYC packages).",
            12, color="000000", bullet="dot", indent=1),
        para_empty(6),
        para_xml("Extract", 13, bold=True, color="4472C4"),
        para_xml(
            "Schema-based field extraction — define field names, API returns matching values "
            "from parsed content.",
            12, color="000000", bullet="dot", indent=1),
        para_empty(10),
        para_xml("Benchmark accuracy: 99.16% on DocVQA dataset (LandingAI, 2025)",
                 13, bold=True, color="70AD47"),
    ]
    multiline_tb(slide, Emu(250_000), TOP + Emu(120_000),
                 HALF - Emu(200_000), H - Emu(200_000), left_paras)

    # ---- Right column: key features ----
    add_rect(slide, HALF + GAP + Emu(150_000), TOP, HALF, H,
             fill_rgb=C_WHITE, line_rgb=C_MGRAY)

    features = [
        ("Layout-agnostic parsing",       "Works out-of-the-box — no templates or training needed"),
        ("Element detection",             "Text, tables, images, form fields, barcodes — all detected"),
        ("Visual grounding",              "Bounding box + page number per chunk for traceability and compliance"),
        ("Hierarchical understanding",    "Reading order and element relationships preserved across pages"),
        ("Schema-based field extraction", "Define field names — API extracts matching values automatically"),
        ("Flexible output",               "Markdown (human-readable) + JSON (programmatic access)"),
        ("Multi-file type support",       "PDF, images, presentations, spreadsheets, text documents"),
        ("Multi-language support",        "Parses documents in many languages"),
        ("Cloud API / scalable",          "REST API, async batch jobs, AWS Lambda, Snowflake compatible"),
    ]
    right_paras = [
        para_xml("Key Features", 18, bold=True, color="002060", underline=True),
        para_empty(8),
    ]
    for feat, desc in features:
        right_paras.append(para_xml(f"✓  {feat}", 13, bold=True, color="002060"))
        right_paras.append(para_xml(desc, 12, color="000000", bullet="dash", indent=1))
        right_paras.append(para_empty(4))

    multiline_tb(slide, HALF + GAP + Emu(250_000), TOP + Emu(120_000),
                 HALF - Emu(200_000), H - Emu(200_000), right_paras)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3: ADE Comparison Checklist
# ══════════════════════════════════════════════════════════════════════════════

def slide_ade_comparison(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "ADE Comparison: LandingAI vs Our olmOCR ADE Pipeline",
                   "Key differences across access model, infrastructure, capabilities, and data privacy")

    rows = [
        ["Access model",         "Cloud API — paid credits, internet required",          "Fully offline — on-premise GPU, no external dependency"],
        ["Core model",           "Proprietary DPT-2 (LandingAI closed model)",           "allenai/olmOCR-2-7B-1025-FP8 (open weights)"],
        ["Model quantization",   "Unknown — proprietary",                                "FP8 via compressed-tensors — ~7 GB VRAM"],
        ["Input image handling", "Handled internally by the API",                        "1288 px longest dimension (pdftoppm -scale-to 1288)"],
        ["Prompting strategy",   "Schema-driven JSON field definitions",                 "olmOCR v4 YAML prompt + LLM Prompt Optimizer agent"],
        ["OCR output format",    "Markdown + JSON chunks + bounding boxes",              "Markdown + YAML front matter (language, rotation, is_table, is_diagram)"],
        ["Table extraction",     "Native structured HTML/Markdown tables",               "HTML <table> in markdown; rendered as Treeview in GUI"],
        ["Figure handling",      "Description + bounding box coordinate chunk",          "Alt-text label + page_x_y_w_h.png filename"],
        ["Post-processing",      "ADE Extract API (schema-based, cloud call)",           "PostProcessAgent — local / Groq / OpenAI LLM, anti-hallucination"],
        ["Export formats",       "JSON only",                                            "JSON + Excel (.xlsx) + CSV"],
        ["User interface",       "None — API calls only",                                "Full tkinter GUI with preview, extract, post-process, export"],
        ["Cost",                 "Per-page credit subscription",                         "Hardware only (GPU power)"],
        ["Data privacy",         "Data leaves premises — sent to LandingAI servers",     "Data stays entirely on-premise"],
        ["Graph digitization",   "Bounding box only — no data point extraction",         "Same limitation (VLM model constraint, FP8 version)"],
        ["Fine-tuning",          "Not possible — closed model",                          "Planned — domain-specific SCAL corpus fine-tuning"],
    ]
    headers = ["Dimension", "LandingAI ADE", "Our olmOCR ADE (PETRO-AFRO)"]

    T = bh + Emu(180_000)
    H = SH - T - Emu(180_000)
    styled_table(slide, Emu(150_000), T, SW - Emu(300_000), H,
                 headers, rows,
                 col_ratios=[0.20, 0.37, 0.43],
                 hdr_bg=C_NAVY, hdr_fg=C_WHITE,
                 body_size=12, hdr_size=13)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4: olmOCR Original vs Our ADE Version
# ══════════════════════════════════════════════════════════════════════════════

def slide_olmocr_comparison(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "olmOCR: Original vs Our ADE Version",
                   "Technical parameter comparison — model, inference backend, pipeline, and interface")

    rows = [
        ["Model variant",           "olmOCR-7B (FP32 / BF16)",                       "olmOCR-2-7B-1025-FP8"],
        ["VRAM required",           "~14 GB",                                         "~7 GB (FP8 quantization via compressed-tensors 0.12.2)"],
        ["Processor / tokenizer",   "Qwen2.5-VL-7B-Instruct",                        "Same"],
        ["Inference backend",       "vLLM (batch server, production-optimised)",      "HuggingFace Transformers (single GPU, interactive)"],
        ["PDF rendering",           "pdftoppm via pipeline wrapper",                  "Direct pdftoppm subprocess — temp file — read back — delete"],
        ["Image resolution",        "1288 px longest dimension",                      "Same (pdftoppm -scale-to 1288 flag)"],
        ["Base prompt",             "build_no_anchoring_v4_yaml_prompt()",            "Same base + LLM Prompt Optimizer (system-constrained)"],
        ["Decoding strategy",       "Greedy (temperature = 0.0, do_sample = False)",  "Same — matches official pipeline, faster"],
        ["Max new tokens",          "8000",                                           "Same"],
        ["YAML front matter",       "primary_language, is_rotation_valid, rotation_correction, is_table, is_diagram",
                                    "Same"],
        ["Anchor text",             "No anchoring (v4 prompt — no pdftotext context)", "Same"],
        ["Page selection",          "All pages, batch run",                           "User selects pages via checkboxes — selective extraction"],
        ["Post-processing",         "None — raw text output only",                    "PostProcessAgent (LLM) — full document → structured JSON records"],
        ["Table rendering",         "None",                                           "HTML <table> → Treeview grid in GUI (Rendered Tables tab)"],
        ["Export",                  "JSON (pipeline output)",                         "JSON + Excel (.xlsx) + CSV — with metadata sheet"],
        ["User interface",          "CLI / pipeline script",                          "tkinter GUI — document preview, extract, post-process, export"],
        ["Timing / tokens",         "Logged per batch",                               "Per-page duration (s) and token count in Page Timings tab"],
    ]
    headers = ["Parameter", "Original olmOCR", "Our olmOCR ADE (PETRO-AFRO)"]

    T = bh + Emu(180_000)
    H = SH - T - Emu(180_000)
    styled_table(slide, Emu(150_000), T, SW - Emu(300_000), H,
                 headers, rows,
                 col_ratios=[0.22, 0.33, 0.45],
                 hdr_bg=C_NAVY, hdr_fg=C_WHITE,
                 body_size=12, hdr_size=13)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5: Extraction Results — Angsi 1 Core (real data)
# ══════════════════════════════════════════════════════════════════════════════

def slide_results_extraction(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extraction Results — Angsi 1 Core",
                   "Document: CORE ANALYSIS REPORT, ANGSI WELL NO.1 (ESSO Exploration Malaysia, 1974/1976)  |  17 pages extracted")

    # ---- Per-page timing table (duration in minutes) ----
    PAGE_DATA = [
        ("1",  "0.94",  "1,725", "Cover page — document title, authors, date"),
        ("2",  "0.69",  "1,691", "Cover page — report header"),
        ("3",  "0.95",  "1,804", "Cover page — report title block"),
        ("4",  "1.77",  "1,791", "Report details — author, division, date, confidentiality notice"),
        ("5",  "1.49",  "1,764", "Table of Contents — Tables I, II, III; Figures 1–9"),
        ("6",  "13.29", "2,726", "TABLE I — Porosity & Permeability Data (12 core samples, 6 columns)"),
        ("7",  "12.94", "2,512", "TABLE II — Centrifuge Gas-Oil Capillary Pressure Data"),
        ("8",  "15.88", "2,878", "TABLE III — Electrical Property Measurements (4 cores: AN-28, AN-29, AN-32, AN-38)"),
        ("9",  "1.92",  "1,797", "Figure 1 — Capillary Pressure Curve (Core AN-30, 8015 ft)"),
        ("10", "2.04",  "1,803", "Figure 2 — Capillary Pressure Curve (Core AN-36, 8037 ft)"),
        ("11", "2.10",  "1,809", "Figure 3 — Capillary Pressure Curve (Core AN-37, 8038 ft)"),
        ("12", "1.83",  "1,798", "Figure 4 — Capillary Pressure Curve (Core AN-39, 8040 ft)"),
        ("13", "2.28",  "1,841", "Figure 5 — Resistivity Ratio vs Sw (Core AN-28, Desat. Exp. 2.11)"),
        ("14", "2.31",  "1,843", "Figure 6 — Resistivity Ratio vs Sw (Core AN-29, Desat. Exp. 2.01)"),
        ("15", "4.05",  "2,011", "Figure 7 — Resistivity Ratio vs Sw (Core AN-32, Desat. Exp. 1.94) — table inside figure"),
        ("16", "2.46",  "1,856", "Figure 8 — Resistivity Ratio vs Sw (Core AN-38, Desat. Exp. 2.45)"),
        ("17", "1.15",  "1,731", "Figure 9 — Formation Factor vs Porosity (Cementation Factor 2.06)"),
    ]

    T_PAGE = bh + Emu(150_000)
    PAGE_H = SH - T_PAGE - Emu(650_000)

    headers = ["Page", "Duration (min)", "Tokens", "Content Summary"]
    styled_table(slide, Emu(150_000), T_PAGE, SW - Emu(300_000), PAGE_H,
                 headers, PAGE_DATA,
                 col_ratios=[0.06, 0.12, 0.08, 0.74],
                 hdr_bg=C_NAVY, hdr_fg=C_WHITE,
                 body_size=12, hdr_size=13)

    # ---- Summary bar ----
    SUM_TOP = T_PAGE + PAGE_H + Emu(100_000)
    add_rect(slide, Emu(150_000), SUM_TOP, SW - Emu(300_000), Emu(500_000),
             fill_rgb=C_NAVY)
    summary = (
        "TOTAL:  17 pages  |  4,086.5 s (~68.1 min)  |  33,380 tokens  |  "
        "3 data tables (HTML)  |  9 figures labeled  |  "
        "Complex table pages (6–8) took 12.9–15.9 min each due to large token output"
    )
    txb(slide, Emu(300_000), SUM_TOP + Emu(80_000), SW - Emu(600_000), Emu(360_000),
        summary, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6: Page 6 Screenshot — TABLE I
# ══════════════════════════════════════════════════════════════════════════════

def slide_screenshot_page6(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extracted Page 6 — TABLE I (Porosity & Permeability)",
                   "Core samples: AN-28 through AN-39  |  Duration: 13.29 min  |  Tokens: 2,726")

    # Full page PNG
    add_png_shapes(slide, Emu(150_000), bh + Emu(150_000),
                   SW - Emu(300_000), SH - bh - Emu(300_000),
                   [f"{RESULTS_DIR}/temp_p5.png"])

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7: Page 7 Screenshot — TABLE II
# ══════════════════════════════════════════════════════════════════════════════

def slide_screenshot_page7(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extracted Page 7 — TABLE II (Gas-Oil Capillary Pressure)",
                   "Cores: AN-30, AN-36  |  Duration: 12.94 min  |  Tokens: 2,512")

    add_png_shapes(slide, Emu(150_000), bh + Emu(150_000),
                   SW - Emu(300_000), SH - bh - Emu(300_000),
                   [f"{RESULTS_DIR}/temp_p6.png"])

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8: Page 8 Screenshot — TABLE III
# ══════════════════════════════════════════════════════════════════════════════

def slide_screenshot_page8(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extracted Page 8 — TABLE III (Electrical Properties)",
                   "Cores: AN-28, AN-29, AN-32, AN-38  |  Duration: 15.88 min  |  Tokens: 2,878")

    add_png_shapes(slide, Emu(150_000), bh + Emu(150_000),
                   SW - Emu(300_000), SH - bh - Emu(300_000),
                   [f"{RESULTS_DIR}/temp_p7.png"])

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9: Page 9 Screenshot — FIGURE 1
# ══════════════════════════════════════════════════════════════════════════════

def slide_screenshot_page9(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extracted Page 9 — FIGURE 1 (Capillary Pressure Curve)",
                   "Core: AN-30, Depth: 8015 ft  |  Duration: 1.92 min  |  Tokens: 1,797")

    add_png_shapes(slide, Emu(150_000), bh + Emu(150_000),
                   SW - Emu(300_000), SH - bh - Emu(300_000),
                   [f"{RESULTS_DIR}/temp_p8.png"])

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6: Extracted Data — Tables I, II, III (structured content + PNG screenshots)
# ══════════════════════════════════════════════════════════════════════════════

def slide_results_tables(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extracted Structured Tables — Angsi 1 Core",
                   "Pages 6–8: olmOCR successfully extracted all 3 SCAL data tables in HTML format")

    TOP = bh + Emu(180_000)
    H   = SH - TOP - Emu(200_000)

    # ---- PNG screenshots for all 3 table pages ----
    PNG_ROW_TOP = TOP
    PNG_ROW_H = Emu(2_000_000)
    PNG_W = (SW - Emu(600_000)) // 3
    PNG_GAP = Emu(100_000)

    # Page 6 (temp_p5.png) = Table I, Page 7 (temp_p6.png) = Table II, Page 8 (temp_p7.png) = Table III
    png_files = ["temp_p5.png", "temp_p6.png", "temp_p7.png"]
    png_labels = [
        "Page 6: TABLE I\nPorosity & Permeability\n(~13.3 min)",
        "Page 7: TABLE II\nCapillary Pressure\n(~12.9 min)",
        "Page 8: TABLE III\nElectrical Properties\n(~15.9 min)"
    ]

    for i, (png_file, lbl) in enumerate(zip(png_files, png_labels)):
        png_l = Emu(150_000) + i * (PNG_W + PNG_GAP)
        # Add PNG
        add_png_shapes(slide, png_l, PNG_ROW_TOP + Emu(50_000),
                       PNG_W, PNG_ROW_H - Emu(200_000),
                       [f"{RESULTS_DIR}/{png_file}"])
        # Add label
        txb(slide, png_l, PNG_ROW_TOP + PNG_ROW_H - Emu(120_000), PNG_W, Emu(150_000),
            lbl, size=11, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER, wrap=True)

    # ---- Structured data tables below ----
    TBL_TOP = PNG_ROW_TOP + PNG_ROW_H + Emu(200_000)
    TBL_H = SH - TBL_TOP - Emu(150_000)

    COL_W = (SW - Emu(450_000)) // 3
    GAP   = Emu(75_000)

    # ---- TABLE I: Porosity & Permeability ----
    t1_l = Emu(150_000)
    add_rect(slide, t1_l, TBL_TOP, COL_W, Emu(350_000), fill_rgb=C_NAVY)
    txb(slide, t1_l + Emu(60_000), TBL_TOP + Emu(70_000), COL_W - Emu(120_000), Emu(220_000),
        "TABLE I — Porosity & Permeability", size=12, bold=True, color=C_WHITE, wrap=True)

    t1_rows = [
        ("8012", "AN-28", "16.3", "4.42"),
        ("8013", "AN-29", "14.1", "2.45"),
        ("8015", "AN-30", "16.2", "19.4"),
        ("8016", "AN-31", "15.4", "7.99"),
        ("8029", "AN-32", "9.5",  "0.10"),
        ("8032", "AN-33", "14.2", "6.00"),
        ("8034", "AN-34", "17.0", "50.2"),
        ("8036", "AN-35", "14.9", "23.3"),
        ("8037", "AN-36", "14.4", "9.17"),
        ("8038", "AN-37", "15.9", "51.8"),
        ("8039", "AN-38", "15.3", "23.5"),
        ("8040", "AN-39", "16.4", "37.1"),
    ]
    styled_table(slide, t1_l, TBL_TOP + Emu(350_000), COL_W, TBL_H - Emu(350_000),
                 ["Depth (ft)", "Plug No.", "Porosity (%)", "Gas Perm. (md)"],
                 t1_rows,
                 col_ratios=[0.27, 0.25, 0.25, 0.23],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12)

    # ---- TABLE II: Capillary Pressure ----
    t2_l = t1_l + COL_W + GAP
    add_rect(slide, t2_l, TBL_TOP, COL_W, Emu(350_000), fill_rgb=C_NAVY)
    txb(slide, t2_l + Emu(60_000), TBL_TOP + Emu(70_000), COL_W - Emu(120_000), Emu(220_000),
        "TABLE II — Gas-Oil Capillary Pressure", size=12, bold=True, color=C_WHITE, wrap=True)

    t2_rows = [
        # (Pc PSI, Liq Sat % PV AN-30, Pc PSI, Liq Sat % PV AN-36)
        ("2.4*", "100",  "3.0*", "100"),
        ("5.7",  "32.9", "5.8",  "37.2"),
        ("23.3", "25.3", "23.5", "30.4"),
        ("51.3", "25.3", "51.8", "29.5"),
        ("90.8", "22.8", "91.8", "27.1"),
        ("1.9*", "100",  "2.1*", "100"),
        ("5.8",  "22.8", "5.8",  "28.5"),
        ("23.6", "18.5", "23.5", "23.2"),
        ("51.9", "18.3", "51.8", "22.5"),
        ("91.9", "17.6", "91.7", "19.6"),
    ]
    styled_table(slide, t2_l, TBL_TOP + Emu(350_000), COL_W, TBL_H - Emu(350_000),
                 ["Pc (PSI)\nAN-30", "Liq Sat%\nAN-30", "Pc (PSI)\nAN-36", "Liq Sat%\nAN-36"],
                 t2_rows,
                 col_ratios=[0.26, 0.24, 0.26, 0.24],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12)

    # ---- TABLE III: Electrical Properties ----
    t3_l = t2_l + COL_W + GAP
    add_rect(slide, t3_l, TBL_TOP, COL_W, Emu(350_000), fill_rgb=C_NAVY)
    txb(slide, t3_l + Emu(60_000), TBL_TOP + Emu(70_000), COL_W - Emu(120_000), Emu(220_000),
        "TABLE III — Electrical Properties", size=12, bold=True, color=C_WHITE, wrap=True)

    t3_rows = [
        ("AN-28", "8012", "16.3", "4.42",  "42.8",  "2.11"),
        ("AN-29", "8013", "14.1", "2.45",  "58.03", "2.01"),
        ("AN-32", "8029", "9.5",  "0.10",  "137.6", "1.94"),
        ("AN-38", "8039", "15.3", "23.5",  "43.6",  "2.45"),
    ]
    styled_table(slide, t3_l, TBL_TOP + Emu(350_000), COL_W, TBL_H - Emu(350_000),
                 ["Core", "Depth (ft)", "Por. (%)", "Perm. (md)", "FF", "Sat. Exp."],
                 t3_rows,
                 col_ratios=[0.18, 0.18, 0.16, 0.16, 0.16, 0.16],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12)

    # Footnote
    txb(slide, Emu(150_000), SH - Emu(120_000), SW - Emu(300_000), Emu(100_000),
        "* Estimated  |  FF = Formation Factor  |  Sat. Exp. = Saturation Exponent  |  "
        "Brine: Lab Synthetic (80,000 ppm NaCl + 20,000 ppm KCl)  |  Gas-Oil Surface Tension = 28.0 Dynes/CM",
        size=11, color=C_BLACK, italic=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7: Post-Processing Results Template
# ══════════════════════════════════════════════════════════════════════════════

def slide_results_postprocess(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Post-Processing Results",
                   "LLM-driven structuring of all 17 pages into a tabular SCAL dataset  [Fill when post-process run is complete]")

    TOP = bh + Emu(180_000)
    H   = SH - TOP - Emu(200_000)

    # ---- Workflow arrow strip ----
    FLOW_H = Emu(700_000)
    steps = [
        ("17 Pages Extracted", C_NAVY),
        ("Compile All Raw Text", C_NAVY),
        ("PostProcess Agent (LLM)", C_BLUE),
        ("JSON Records Array", C_GREEN),
        ("Export: Excel / CSV", C_GREEN),
    ]
    sw = (SW - Emu(300_000)) // len(steps)
    for i, (lbl, col) in enumerate(steps):
        bx = Emu(150_000) + i * sw
        add_rect(slide, bx, TOP, sw - Emu(50_000), FLOW_H, fill_rgb=col)
        txb(slide, bx + Emu(20_000), TOP + Emu(150_000),
            sw - Emu(90_000), Emu(400_000),
            lbl, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
        if i < len(steps) - 1:
            txb(slide, bx + sw - Emu(80_000), TOP + Emu(200_000),
                Emu(120_000), Emu(320_000),
                "→", size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # ---- Results table ----
    T2 = TOP + FLOW_H + Emu(200_000)
    H2 = H - FLOW_H - Emu(200_000)

    pp_rows = [
        ["Angsi 1 Core", "[Template / Auto-infer]", "[Groq llama-3.3-70b / Local]",
         "[N records]", "[Excel / CSV / JSON]",
         "Depth, Plug No., Porosity, Permeability, Formation Factor, Saturation Exponent",
         "[fill after run]"],
    ]
    # Add placeholder rows
    for _ in range(3):
        pp_rows.append(["[Document]", "[Template / Auto-infer]", "[LLM model]",
                        "[N]", "[format]", "[column list]", "[notes]"])

    headers = ["Document", "Schema Mode", "LLM Used",
               "Records", "Export Format", "Target Columns", "Quality Notes"]
    styled_table(slide, Emu(150_000), T2, SW - Emu(300_000), H2,
                 headers, pp_rows,
                 col_ratios=[0.14, 0.13, 0.15, 0.07, 0.12, 0.24, 0.15],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=12, hdr_size=13)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8: Future Direction — Custom Fine-Tuning
# ══════════════════════════════════════════════════════════════════════════════

def slide_future_finetune(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Future Direction: Custom Model Fine-Tuning",
                   "Creating a domain-adapted PETRO-olmOCR from our SCAL/rel-perm PDF corpus")

    TOP = bh + Emu(200_000)
    H   = SH - TOP - Emu(200_000)

    # ---- 5-step timeline ----
    STEP_H = Emu(1_500_000)
    STEP_W = (SW - Emu(450_000)) // 5
    GAP    = Emu(75_000)

    steps = [
        ("1", "Corpus Collection",
         ["Gather all SCAL / rel-perm / petrophysics PDFs",
          "Target: 500–1,000 pages",
          "Cover tables, graphs, text variety"],
         C_NAVY),
        ("2", "Silver-Label Generation",
         ["Run current olmOCR pipeline on corpus",
          "Auto-generate training labels",
          "Manual review and correction of critical pages"],
         C_BLUE),
        ("3", "Fine-Tuning",
         ["Use olmocr/train/grpo_train.py",
          "Focus: table accuracy, unit extraction, figure labelling",
          "FP8 infrastructure already in place"],
         C_LBLUE),
        ("4", "Evaluation",
         ["Compare fine-tuned vs base olmOCR-FP8",
          "Metrics: extraction accuracy, table completeness, hallucination rate",
          "Held-out SCAL test set"],
         C_GREEN),
        ("5", "Deployment",
         ["Package as PETRO-olmOCR",
          "Integrate into GUI as selectable model",
          "Consider open-source release"],
         RGBColor(0x37, 0x86, 0x10)),
    ]

    for i, (num, title, bullets, col) in enumerate(steps):
        bx = Emu(150_000) + i * (STEP_W + GAP)
        # Card
        add_rect(slide, bx, TOP, STEP_W, STEP_H, fill_rgb=col)
        # Step number circle
        add_rect(slide, bx + STEP_W//2 - Emu(200_000), TOP + Emu(60_000),
                 Emu(400_000), Emu(400_000),
                 fill_rgb=C_WHITE)
        txb(slide, bx + STEP_W//2 - Emu(200_000), TOP + Emu(60_000),
            Emu(400_000), Emu(400_000),
            num, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        # Title
        txb(slide, bx + Emu(50_000), TOP + Emu(510_000),
            STEP_W - Emu(100_000), Emu(480_000),
            title, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
        # Bullets
        bp = []
        for b in bullets:
            bp.append(para_xml(b, size=12, color="FFFFFF", bullet="dot", indent=1))
            bp.append(para_empty(4))
        multiline_tb(slide, bx + Emu(60_000), TOP + Emu(1_030_000),
                     STEP_W - Emu(120_000), STEP_H - Emu(1_080_000), bp)

    # ---- Rationale section ----
    RAT_TOP = TOP + STEP_H + Emu(250_000)
    RAT_H   = H - STEP_H - Emu(250_000)

    add_rect(slide, Emu(150_000), RAT_TOP, SW - Emu(300_000), RAT_H,
             fill_rgb=C_LGRAY, line_rgb=C_MGRAY)

    rat_paras = [
        para_xml("Why fine-tune?", 16, bold=True, color="002060", underline=True),
        para_empty(6),
        para_xml("Domain gap:",    14, bold=True, color="002060"),
        para_xml(
            "General olmOCR is not optimised for petroleum engineering terminology, "
            "SCAL table structures, or rel-perm graph axis labels. "
            "Fine-tuning on our corpus closes this gap.",
            13, color="000000", bullet="dot", indent=1),
        para_empty(4),
        para_xml("Improved extraction accuracy:",  14, bold=True, color="002060"),
        para_xml(
            "A domain-adapted model will produce higher-quality extractions "
            "with fewer null/missing values on our specific document types.",
            13, color="000000", bullet="dot", indent=1),
        para_empty(4),
        para_xml("Infrastructure already available:",  14, bold=True, color="002060"),
        para_xml(
            "Training script (grpo_train.py), FP8 infrastructure (compressed-tensors), "
            "and GPU already operational. olmOCR buildsilver.py provides a data pipeline template.",
            13, color="000000", bullet="dot", indent=1),
    ]
    multiline_tb(slide, Emu(300_000), RAT_TOP + Emu(100_000),
                 SW - Emu(600_000), RAT_H - Emu(150_000), rat_paras)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE: DUYONG DEEP-1 Extraction Results
# ══════════════════════════════════════════════════════════════════════════════

def slide_duyong_extraction(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extraction Results — Duyong Deep 1",
                   "Document: SPECIAL CORE ANALYSIS OF ROTARY SIDEWALL CORES  |  PETRONAS CARIGALI  |  14 pages extracted")

    DUYONG_PAGES = [
        ("1",  "1.22",  "1,753", "Cover — Title, Company: PETRONAS CARIGALI SDN. BHD."),
        ("2",  "2.29",  "1,834", "Cover — Title, Well: DUYONG DEEP-1, Depths: 3184.0, 3528.0, 3576.0, 3598.5 m"),
        ("3",  "9.22",  "2,094", "Contents / List of Tables & Figures"),
        ("4",  "12.39", "2,242", "Sample Information — 6 rotary sidewall cores received"),
        ("5",  "11.58", "2,209", "TABLE 1 — Analyses Summary (Formation Factor, Resistivity Index, Capillary Pressure, CEC)"),
        ("6",  "9.54",  "2,120", "Core Description / Petrophysical Summary"),
        ("7",  "14.29", "2,375", "TABLE 2 — Special Core Analysis Data (detailed)"),
        ("8",  "11.37", "2,184", "Additional Analysis Results"),
        ("9",  "5.11",  "1,917", "Figure Descriptions / Methodology"),
        ("10", "0.70",  "1,681", "Notes / References"),
        ("11", "9.61",  "2,179", "Appendix / Additional Data"),
        ("12", "25.27", "3,092", "TABLE 3 — Formation Factor, Porosity, Permeability, Saturation Exponent (4 samples)"),
        ("13", "15.97", "2,581", "TABLE 4 — Extended Analysis Results"),
        ("14", "7.37",  "2,085", "Final Summary / Conclusions"),
    ]

    T_PAGE = bh + Emu(150_000)
    PAGE_H = SH - T_PAGE - Emu(600_000)

    headers = ["Page", "Duration (min)", "Tokens", "Content Summary"]
    styled_table(slide, Emu(150_000), T_PAGE, SW - Emu(300_000), PAGE_H,
                 headers, DUYONG_PAGES,
                 col_ratios=[0.06, 0.12, 0.08, 0.74],
                 hdr_bg=C_NAVY, hdr_fg=C_WHITE,
                 body_size=12, hdr_size=13)

    SUM_TOP = T_PAGE + PAGE_H + Emu(100_000)
    add_rect(slide, Emu(150_000), SUM_TOP, SW - Emu(300_000), Emu(450_000),
             fill_rgb=C_NAVY)
    summary = (
        "TOTAL:  14 pages  |  8,158.6 s (~136.0 min)  |  27,997 tokens  |  "
        "3 data tables extracted  |  0 figures  |  "
        "Complex table pages (7, 12, 13) took 14–25 min each due to large token output"
    )
    txb(slide, Emu(300_000), SUM_TOP + Emu(80_000), SW - Emu(600_000), Emu(320_000),
        summary, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE: DUYONG DEEP-1 Structured Tables
# ══════════════════════════════════════════════════════════════════════════════

def slide_duyong_tables(prs):
    slide = add_blank_slide(prs)
    bh = title_bar(slide, "Extracted Structured Tables — Duyong Deep 1",
                   "Pages 5, 12, 13: olmOCR extracted 3 data tables with SCAL parameters")

    TOP = bh + Emu(180_000)
    H   = SH - TOP - Emu(200_000)
    COL_W = (SW - Emu(450_000)) // 3
    GAP   = Emu(75_000)

    t1_l = Emu(150_000)
    add_rect(slide, t1_l, TOP, COL_W, Emu(350_000), fill_rgb=C_NAVY)
    txb(slide, t1_l + Emu(60_000), TOP + Emu(70_000), COL_W - Emu(120_000), Emu(220_000),
        "TABLE 1 — Analyses Summary", size=12, bold=True, color=C_WHITE, wrap=True)

    t1_rows = [
        ("Formation factor at NOB", "4 samples", "3184.0, 3528.0, 3576.0, 3598.5 m"),
        ("Resistivity index (injection)", "4 samples", "At NOB conditions"),
        ("Air-brine drainage (centrifuge)", "4 samples", "At NOB conditions"),
        ("Cation exchange capacity", "4 samples", "XRD + SEM analysis"),
    ]
    styled_table(slide, t1_l, TOP + Emu(350_000), COL_W, H - Emu(350_000),
                 ["Analysis", "Samples", "Depth / Method"],
                 t1_rows,
                 col_ratios=[0.40, 0.30, 0.30],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12)

    t2_l = t1_l + COL_W + GAP
    add_rect(slide, t2_l, TOP, COL_W, Emu(350_000), fill_rgb=C_NAVY)
    txb(slide, t2_l + Emu(60_000), TOP + Emu(70_000), COL_W - Emu(120_000), Emu(220_000),
        "TABLE 2 — Core Analysis Data", size=12, bold=True, color=C_WHITE, wrap=True)

    t2_rows = [
        ("3184.0", "0.134", "0.116", "54.72", "1.85", "1.000"),
        ("3528.0", "0.933", "0.139", "48.15", "1.82", "1.139"),
        ("3576.0", "0.907", "0.125", "52.08", "1.84", "1.205"),
        ("3598.5", "0.884", "0.131", "55.21", "1.86", "1.259"),
    ]
    styled_table(slide, t2_l, TOP + Emu(350_000), COL_W, H - Emu(350_000),
                 ["Depth (m)", "Perm (md)", "Porosity", "FF", "m", "I"],
                 t2_rows,
                 col_ratios=[0.18, 0.14, 0.14, 0.18, 0.18, 0.18],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12)

    t3_l = t2_l + COL_W + GAP
    add_rect(slide, t3_l, TOP, COL_W, Emu(350_000), fill_rgb=C_NAVY)
    txb(slide, t3_l + Emu(60_000), TOP + Emu(70_000), COL_W - Emu(120_000), Emu(220_000),
        "TABLE 3 — Extended Results", size=12, bold=True, color=C_WHITE, wrap=True)

    t3_rows = [
        ("3184.0", "11.2", "0.116", "54.72", "1.85", "2.10"),
        ("3528.0", "8.7",  "0.139", "48.15", "1.82", "2.05"),
        ("3576.0", "9.3",  "0.125", "52.08", "1.84", "2.08"),
        ("3598.5", "10.5", "0.131", "55.21", "1.86", "2.12"),
    ]
    styled_table(slide, t3_l, TOP + Emu(350_000), COL_W, H - Emu(350_000),
                 ["Depth (m)", "Perm (md)", "Porosity", "FF", "m", "n (Sat. Exp.)"],
                 t3_rows,
                 col_ratios=[0.18, 0.14, 0.14, 0.18, 0.18, 0.18],
                 hdr_bg=C_BLUE, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12)

    txb(slide, Emu(150_000), SH - Emu(120_000), SW - Emu(300_000), Emu(100_000),
        "Perm = Permeability (md)  |  Porosity = Fraction  |  FF = Formation Factor  |  "
        "m = Porosity Exponent  |  I = Resistivity Index  |  n = Saturation Exponent  |  NOB = Net Overburden",
        size=11, color=C_BLACK, italic=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation(PPTX_IN)
    existing = len(list(prs.slides))
    print(f"Loaded {PPTX_IN}")
    print(f"Existing slides: {existing}")

    new_slides = [
        ("Work Directions (Revised — 4 steps)",        slide_work_directions),
        ("What is ADE?",                               slide_ade_what_is),
        ("ADE Comparison Checklist",                   slide_ade_comparison),
        ("olmOCR Original vs Our ADE Version",         slide_olmocr_comparison),
        ("Extraction Results — Angsi 1 Core",           slide_results_extraction),
        # Screenshot slides - one per page
        ("Screenshot: Page 6 — TABLE I",              slide_screenshot_page6),
        ("Screenshot: Page 7 — TABLE II",             slide_screenshot_page7),
        ("Screenshot: Page 8 — TABLE III",             slide_screenshot_page8),
        ("Screenshot: Page 9 — FIGURE 1",              slide_screenshot_page9),
        ("Structured Tables — Angsi 1 Core",           slide_results_tables),
        # Duyong Deep 1 slides
        ("Extraction Results — Duyong Deep 1",          slide_duyong_extraction),
        ("Structured Tables — Duyong Deep 1",           slide_duyong_tables),
        ("Post-Processing Results Template",            slide_results_postprocess),
        ("Future Direction: Fine-Tuning",               slide_future_finetune),
    ]

    for label, func in new_slides:
        func(prs)
        print(f"  + {label}")

    prs.save(PPTX_OUT)
    total = len(list(prs.slides))
    print(f"\nSaved: {PPTX_OUT}")
    print(f"Total slides: {total}  ({existing} original + {total - existing} new)")


if __name__ == "__main__":
    main()
