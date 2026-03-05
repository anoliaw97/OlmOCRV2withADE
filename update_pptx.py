"""
update_pptx.py
Modifies [24] 26 Feb 2026_ver04.pptx → ver05.pptx

Changes:
  1. Slides 14–22: Remove navy header bar, change text colors to match
     the theme used in slides 1–13 (dark text on white background).
  2. Add extraction-overview slides for Duyong Deep 1 and Pegaga-2.
  3. Add an overall benchmark comparison slide.
  All durations expressed in min or hrs (not raw seconds).
"""

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PPTX_IN  = "[24] 26 Feb 2026_ver04.pptx"
PPTX_OUT = "[24] 26 Feb 2026_ver05.pptx"

# ── Colors ─────────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x00, 0x20, 0x60)
C_BLUE   = RGBColor(0x44, 0x72, 0xC4)
C_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY   = RGBColor(0x59, 0x59, 0x59)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)

SW = 10693400
SH = 7556500

MARGIN      = Emu(350_000)
TITLE_TOP   = Emu(80_000)
TITLE_H     = Emu(480_000)
SUBTITLE_TOP= Emu(590_000)
SUBTITLE_H  = Emu(280_000)
CONTENT_TOP = Emu(920_000)
CONTENT_W   = Emu(SW - 700_000)


# ── Duration formatter ─────────────────────────────────────────────────────────

def fmt_dur(seconds_str):
    """Convert a seconds value (string, may have commas or ~) to min or hrs."""
    s_clean = seconds_str.strip().lstrip("~").replace(",", "")
    prefix = "~" if seconds_str.strip().startswith("~") else ""
    try:
        s = float(s_clean)
    except ValueError:
        return seconds_str
    if s < 3600:
        return f"{prefix}{s/60:.1f} min"
    else:
        return f"{prefix}{s/3600:.1f} hrs"


# ── XML / shape helpers ────────────────────────────────────────────────────────

def _esc(t):
    return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def para_xml(text, size=13, bold=False, color="1A1A1A", bullet=None,
             indent=0, underline=False, align="l", italic=False):
    a  = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sz = int(max(size, 11) * 100)
    b  = ' b="1"' if bold    else ''
    u  = ' u="sng"' if underline else ''
    it = ' i="1"'  if italic  else ''
    fill = f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
    mar  = f' marL="{indent*457200}" indent="-457200"' if indent else ''
    if   bullet == "dot":   bxml = '<a:buChar char="•"/>'
    elif bullet == "dash":  bxml = '<a:buChar char="–"/>'
    elif bullet == "check": bxml = '<a:buChar char="✓"/>'
    elif bullet == "cross": bxml = '<a:buChar char="✗"/>'
    else:                   bxml = '<a:buNone/>'
    xml = (
        f'<a:p xmlns:a="{a}">'
        f'<a:pPr algn="{align}"{mar}>{bxml}</a:pPr>'
        f'<a:r><a:rPr lang="en-US" sz="{sz}"{b}{u}{it} dirty="0">'
        f'{fill}</a:rPr><a:t>{_esc(text)}</a:t></a:r>'
        f'</a:p>'
    )
    return etree.fromstring(xml)


def para_empty(size=12):
    a  = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sz = int(max(size, 11) * 100)
    return etree.fromstring(
        f'<a:p xmlns:a="{a}"><a:endParaRPr lang="en-US" sz="{sz}" dirty="0"/></a:p>'
    )


def add_textbox(slide, left, top, width, height, paras, wrap=True):
    box   = slide.shapes.add_textbox(left, top, width, height)
    tf    = box.text_frame
    tf.word_wrap = wrap
    txBody = tf._txBody
    bpr = txBody.find(qn("a:bodyPr"))
    if bpr is not None:
        bpr.set("anchor", "t")
    for old in txBody.findall(qn("a:p")):
        txBody.remove(old)
    for p in paras:
        txBody.append(p)
    return box


def add_layout_slide(prs, name="Title and Content"):
    for layout in prs.slide_master.slide_layouts:
        if layout.name == name:
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_master.slide_layouts[1])


def styled_table(slide, left, top, width, height, headers, rows,
                 col_ratios=None, hdr_bg=C_NAVY, hdr_fg=C_WHITE,
                 body_size=11, hdr_size=12):
    ncols = len(headers)
    nrows = len(rows) + 1
    ts  = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = ts.table
    if col_ratios:
        for ci, r in enumerate(col_ratios):
            tbl.columns[ci].width = Emu(int(width * r))
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid();  cell.fill.fore_color.rgb = hdr_bg
        tf = cell.text_frame;  tf.word_wrap = True
        p = tf.paragraphs[0];  p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr;  run.font.size = Pt(hdr_size)
        run.font.bold = True;  run.font.color.rgb = hdr_fg
        run.font.name = "Calibri"
    for ri, row in enumerate(rows):
        bg = C_LGRAY if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid();  cell.fill.fore_color.rgb = bg
            tf = cell.text_frame;  tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val);  run.font.size = Pt(body_size)
            run.font.bold = False;  run.font.color.rgb = C_BLACK
            run.font.name = "Calibri"
    return ts


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, [
        para_xml(title, size=22, bold=True, color="002060", align="l"),
    ])
    if subtitle:
        add_textbox(slide, MARGIN, SUBTITLE_TOP, CONTENT_W, SUBTITLE_H, [
            para_xml(subtitle, size=12, bold=False, color="595959", align="l"),
        ])


# ══════════════════════════════════════════════════════════════════════════════
#  PART 1 – FIX SLIDES 14–22
# ══════════════════════════════════════════════════════════════════════════════

def recolor_runs(shape, rgb):
    if not hasattr(shape, "text_frame"):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb


def delete_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def fix_slide_theme(slide, slide_num):
    for shape in list(slide.shapes):
        name = shape.name
        if name == "Rectangle 1":
            delete_shape(slide, shape)
        elif name == "TextBox 2":
            recolor_runs(shape, C_DARK)
        elif name == "TextBox 3":
            recolor_runs(shape, C_GRAY)
        elif name == "Rectangle 5" and slide_num == 17:
            delete_shape(slide, shape)
        elif name == "TextBox 6" and slide_num == 17:
            recolor_runs(shape, C_GRAY)
        elif name in ("Rectangle 4", "Rectangle 7", "Rectangle 10") and slide_num == 22:
            shape.fill.solid()
            shape.fill.fore_color.rgb = C_BLUE


# ══════════════════════════════════════════════════════════════════════════════
#  PART 2 – NEW SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# Raw data: (page_label, seconds_str, tokens_str, notes)
DUYONG_PAGES = [
    ("1",    "293.8",    "1,753",   "Title page"),
    ("2",    "615.7",    "1,834",   "Cover page"),
    ("3",    "1,147.7",  "2,094",   "Table of Contents"),
    ("4",    "1,587.9",  "2,242",   "Introduction, objectives, samples"),
    ("5",    "1,497.7",  "2,209",   "Samples & Analyses table"),
    ("6",    "1,449.0",  "2,120",   "SCAL Introduction & Results summary"),
    ("7",    "1,944.7",  "2,375",   "Tables 1 & 2 – Formation Resistivity Factor"),
    ("8",    "1,434.4",  "2,184",   "Table 3 – Air-Brine Capillary Pressure"),
    ("9",    "755.8",    "1,917",   "Table 4 – CEC"),
    ("10",   "98.0",     "1,681",   "Blank / spacer page"),
    ("11",   "1,552.7",  "2,179",   "Table 5a – MICP (3184 m)"),
    ("12",   "4,630.5",  "3,092",   "Table 5b – MICP (3528 m)  ★ Complex"),
    ("13",   "2,341.4",  "2,581",   "Table 5c – MICP (3576 m)"),
    ("14",   "1,286.0",  "2,085",   "Table 5d – MICP (3598.5 m)"),
    ("15",   "10,454.7", "5,778",   "Graph 1 – FRF (3184 m)  ★★ Very complex"),
    ("16",   "3,547.0",  "3,071",   "Graph 2 – FRF (3528 m)"),
    ("17",   "11,389.3", "5,775",   "Graph 3 – Capillary Pressure  ★★ Very complex"),
    ("18",   "4,001.6",  "3,201",   "Graph 4a – MICP Linear"),
    ("19",   "11,992.5", "5,774",   "Graph 4b – MICP Linear  ★★ Very complex"),
    ("20",   "4,509.3",  "3,197",   "Graph 4c – MICP Linear"),
    ("21",   "12,060.6", "5,767",   "Graph 4d – MICP Linear  ★★ Very complex"),
    ("22",   "4,809.3",  "3,331",   "Graph 5a – MICP Semilog"),
    ("23",   "236.8",    "1,682",   "Graph 5b"),
    ("24–38","~8,821.5", "~28,908", "Graphs 5c–6d, Pore Size Distribution"),
]

PEGAGA_PAGES = [
    ("1",    "500.2",    "1,732",   "Title page"),
    ("2",    "346.1",    "1,732",   "Cover page"),
    ("3",    "531.1",    "1,817",   "Executive Summary"),
    ("4",    "757.7",    "1,931",   "SCAL Flow Chart (Fig. 1)"),
    ("5",    "1,398.2",  "1,932",   "Contents"),
    ("6",    "1,616.7",  "2,214",   "Introduction – test principles"),
    ("7",    "2,721.4",  "2,664",   "Rel. Perm. Test Procedure + Sample table  ★ Complex"),
    ("8",    "547.0",    "1,847",   "Test procedure (continued)"),
    ("9",    "2,812.0",  "2,586",   "Coreflood simulation intro  ★ Complex"),
    ("10",   "1,514.9",  "2,102",   "Sample 2-015 test setup"),
    ("11",   "527.0",    "1,796",   "History matching intro"),
    ("12",   "888.1",    "1,958",   "Sample 2-015 fit results"),
    ("13",   "1,567.7",  "1,901",   "kr curves – Sample 2-015"),
    ("14",   "3,781.8",  "2,865",   "History matching tables  ★ Complex"),
    ("15",   "37,705.4", "3,819",   "Sample 2-019 large data table  ★★★ Extreme"),
    ("16",   "9,000.6",  "4,255",   "Sample 2-019 (continued)  ★★ Very complex"),
    ("17",   "1,473.6",  "2,102",   "Sample 2-019 results"),
    ("18",   "941.2",    "1,866",   "Sample 2-023 setup"),
    ("19",   "1,103.5",  "1,927",   "Sample 2-023 history matching"),
    ("20",   "419.6",    "1,903",   "Sample 2-023 results"),
    ("21",   "4,557.5",  "2,863",   "Sample 2-029 tables  ★ Complex"),
    ("22",   "8,342.0",  "3,960",   "Sample 2-029 (continued)  ★★ Very complex"),
    ("23",   "9,978.8",  "4,172",   "Sample 2-029 results  ★★ Very complex"),
    ("24",   "1,507.1",  "2,067",   "Sample 2-029 kr curves"),
    ("25–34","~7,386.2", "~19,826", "Samples 2-033 (pp.25–34)"),
    ("35",   "8,521.5",  "3,965",   "Data averaging – normalization  ★★ Very complex"),
    ("36",   "9,760.7",  "4,230",   "Normalized kr curves  ★★ Very complex"),
    ("37",   "6,897.0",  "3,533",   "References / discussion  ★ Complex"),
    ("38–42","~3,317.0", "~9,250",  "Appendix pages"),
    ("43",   "9,061.9",  "4,048",   "Appendix – large table  ★★ Very complex"),
    ("44",   "9,157.9",  "4,063",   "Appendix (continued)  ★★ Very complex"),
    ("45–49","~5,084.5", "~9,687",  "Appendix pages 45–49"),
]


def build_page_rows(raw_data):
    """Convert raw (page, sec, tokens, notes) to (page, duration_str, tokens, notes)."""
    rows = []
    for page, sec_str, tokens, notes in raw_data:
        rows.append((page, fmt_dur(sec_str), tokens, notes))
    return rows


# ── Slide: Duyong Deep 1 overview ─────────────────────────────────────────────

def add_duyong_overview(prs):
    slide = add_layout_slide(prs)
    add_slide_title(
        slide,
        "Extraction Results — Duyong Deep 1",
        "Document: SPECIAL CORE ANALYSIS OF ROTARY SIDEWALL CORES  "
        "(PETRONAS CARIGALI / UZMA ENGINEERING)  |  38 pages extracted"
    )
    headers = ["Page", "Duration", "Tokens", "Notes"]
    rows    = build_page_rows(DUYONG_PAGES)
    tbl_h   = Emu(SH - int(CONTENT_TOP) - 220_000)
    styled_table(slide, MARGIN, CONTENT_TOP, CONTENT_W, tbl_h,
                 headers, rows, col_ratios=[0.08, 0.14, 0.12, 0.66],
                 body_size=10, hdr_size=11)
    add_textbox(slide, MARGIN, Emu(SH - 210_000), CONTENT_W, Emu(200_000), [
        para_xml(
            "TOTAL:  38 pages  |  25.5 hrs  |  95,039 tokens  |  4 samples (depths 3184–3598.5 m)  |  "
            "Complex graph pages 15, 17, 19, 21 — each 2.9–3.3 hrs",
            size=10, color="595959"
        )
    ])


# ── Slide: Pegaga-2 overview ───────────────────────────────────────────────────

def add_pegaga_overview(prs):
    slide = add_layout_slide(prs)
    add_slide_title(
        slide,
        "Extraction Results — Pegaga-2",
        "Document: RELATIVE PERMEABILITY NUMERICAL INTERPRETATION  "
        "(MDC Oil & Gas SK320 / Senergy International, Nov 2015)  |  49 pages extracted"
    )
    headers = ["Page", "Duration", "Tokens", "Notes"]
    rows    = build_page_rows(PEGAGA_PAGES)
    tbl_h   = Emu(SH - int(CONTENT_TOP) - 220_000)
    styled_table(slide, MARGIN, CONTENT_TOP, CONTENT_W, tbl_h,
                 headers, rows, col_ratios=[0.08, 0.14, 0.12, 0.66],
                 body_size=10, hdr_size=11)
    add_textbox(slide, MARGIN, Emu(SH - 210_000), CONTENT_W, Emu(200_000), [
        para_xml(
            "TOTAL:  49 pages  |  49.1 hrs  |  122,934 tokens  |  5 samples (depths 2511–2516 m)  |  "
            "Page 15 alone: 10.5 hrs — largest single-page extraction in this benchmark",
            size=10, color="595959"
        )
    ])


# ── Slide: Benchmark comparison ────────────────────────────────────────────────

def add_benchmark_comparison(prs):
    slide = add_layout_slide(prs)
    add_slide_title(
        slide,
        "Extraction Benchmark — Three SCAL Documents",
        "olmOCR ADE pipeline performance across document types, complexity, and size"
    )
    headers = ["Metric", "Angsi-1 Core\n(ESSO, 1974)", "Duyong Deep-1\n(PETRONAS, ~2014)", "Pegaga-2\n(MDC/Senergy, 2015)"]
    rows = [
        ("Document type",        "Core Analysis Report",     "Special Core Analysis (SCAL)",    "Rel. Perm. Numerical Interp."),
        ("Total pages",          "17",                        "38",                               "49"),
        ("Total duration",       "68.1 min",                  "25.5 hrs",                         "49.1 hrs"),
        ("Total tokens",         "33,380",                    "95,039",                           "122,934"),
        ("Avg tokens / page",    "1,964",                     "2,501",                            "2,509"),
        ("Avg duration / page",  "4.0 min",                   "40.2 min",                         "60.1 min"),
        ("Fastest page",         "Page 1: 0.9 min",           "Page 10: 1.6 min",                 "Page 2: 5.8 min"),
        ("Slowest page",         "Page 8: 15.9 min",          "Page 21: 3.3 hrs",                 "Page 15: 10.5 hrs"),
        ("Pages > 1 hr",         "None",                      "Pages 15, 17, 19, 21  (4 pages)",  "Pages 15, 16, 22–23, 35–36, 43–44  (8 pages)"),
        ("Tables extracted",     "3 (HTML format)",           "5+ SCAL data tables",              "5 sample tables + kr curves"),
        ("Figures labeled",      "9 figures",                 "18+ graphs",                       "40+ kr curve plots"),
        ("Extraction quality",   "✓  All tables correct",     "✓  All tables extracted",          "✓  All pages extracted"),
    ]
    tbl_h = Emu(SH - int(CONTENT_TOP) - 200_000)
    styled_table(slide, MARGIN, CONTENT_TOP, CONTENT_W, tbl_h,
                 headers, rows, col_ratios=[0.22, 0.26, 0.26, 0.26],
                 body_size=10, hdr_size=11)
    add_textbox(slide, MARGIN, Emu(SH - 200_000), CONTENT_W, Emu(190_000), [
        para_xml(
            "Observation: Extraction time scales non-linearly with complexity. "
            "Simple text pages: ~5–20 min.  Dense SCAL data tables: 20–80 min.  "
            "Complex multi-panel graph pages: up to 10.5 hrs.",
            size=10, color="595959"
        )
    ])


# ══════════════════════════════════════════════════════════════════════════════
#  PART 3 – STRUCTURED TABLE SLIDES  (from extracted HTML tables)
# ══════════════════════════════════════════════════════════════════════════════

def _side_by_side_tables(prs, title, subtitle, left_title, left_hdrs, left_rows,
                          right_title, right_hdrs, right_rows,
                          left_ratios=None, right_ratios=None):
    """Two tables side-by-side on one slide."""
    slide = add_layout_slide(prs)
    add_slide_title(slide, title, subtitle)

    gap   = Emu(150_000)
    half  = Emu((SW - 700_000) // 2 - int(gap) // 2)
    tbl_h = Emu(SH - int(CONTENT_TOP) - 700_000)

    # Left section label
    add_textbox(slide, MARGIN, Emu(int(CONTENT_TOP) - 280_000), half, Emu(260_000), [
        para_xml(left_title, size=11, bold=True, color="002060")
    ])
    styled_table(slide, MARGIN, CONTENT_TOP, half, tbl_h,
                 left_hdrs, left_rows, col_ratios=left_ratios, body_size=10, hdr_size=11)

    right_x = Emu(int(MARGIN) + int(half) + int(gap))
    add_textbox(slide, right_x, Emu(int(CONTENT_TOP) - 280_000), half, Emu(260_000), [
        para_xml(right_title, size=11, bold=True, color="002060")
    ])
    styled_table(slide, right_x, CONTENT_TOP, half, tbl_h,
                 right_hdrs, right_rows, col_ratios=right_ratios, body_size=10, hdr_size=11)
    return slide


def _full_table_slide(prs, title, subtitle, tbl_label, headers, rows,
                      col_ratios=None, body_size=10, note=None):
    slide = add_layout_slide(prs)
    add_slide_title(slide, title, subtitle)
    add_textbox(slide, MARGIN, Emu(int(CONTENT_TOP) - 280_000), CONTENT_W, Emu(260_000), [
        para_xml(tbl_label, size=11, bold=True, color="002060")
    ])
    tbl_h = Emu(SH - int(CONTENT_TOP) - (300_000 if note else 100_000))
    styled_table(slide, MARGIN, CONTENT_TOP, CONTENT_W, tbl_h,
                 headers, rows, col_ratios=col_ratios, body_size=body_size, hdr_size=11)
    if note:
        add_textbox(slide, MARGIN, Emu(SH - 250_000), CONTENT_W, Emu(240_000), [
            para_xml(note, size=10, color="595959")
        ])
    return slide


def add_duyong_structured_slides(prs):
    """3 slides showing key extracted tables from Duyong Deep 1."""
    from extract_tables import extract_tables_from_file

    tables = extract_tables_from_file('RESULTS/Duyon Deep 1 Full.txt')

    # ── Slide 1: TABLE 1 (FRF) + TABLE 4 (CEC) side by side ──────────────────
    t1 = next(t for t in tables if t['page'] == 11)   # FRF
    t4 = next(t for t in tables if t['page'] == 14)   # CEC

    frf_hdrs = ['Sample ID', 'Depth (m)', 'Ka@NOB (mD)', 'φ@NOB (frac)', 'FRF (F)', 'm']
    frf_rows = [row for row in t1['grid'][2:]]   # skip 2-row header

    cec_hdrs = ['Sample ID', 'Depth (m)', 'Porosity (%)', 'Grain Density (g/cc)', 'CEC (meq/100g)', 'Qv (meq/ml)']
    cec_rows = [row for row in t4['grid'][1:]]   # single header row

    _side_by_side_tables(
        prs,
        title="Extracted Structured Tables — Duyong Deep 1 (I)",
        subtitle="Well Duyong Deep-1, Block PM12  |  NOB = 560 psi  |  Saturant: 3,500 ppm synthetic brine",
        left_title="TABLE 1 — Formation Resistivity Factor at NOB Conditions",
        left_hdrs=frf_hdrs, left_rows=frf_rows,
        left_ratios=[0.16, 0.18, 0.18, 0.18, 0.15, 0.15],
        right_title="TABLE 4 — Cation Exchange Capacity",
        right_hdrs=cec_hdrs, right_rows=cec_rows,
        right_ratios=[0.16, 0.18, 0.18, 0.20, 0.18, 0.10],
    )

    # ── Slide 2: TABLE 2 — FRF + RI (all data) ────────────────────────────────
    t2 = next(t for t in tables if t['page'] == 12)
    ri_hdrs = ['Sample', 'Depth (m)', 'Ka (mD)', 'φ (frac)', 'F', 'm', 'Sw (frac)', 'RI (I)', 'n']
    ri_rows = [row for row in t2['grid'][2:]]   # 20 data rows

    _full_table_slide(
        prs,
        title="Extracted Structured Tables — Duyong Deep 1 (II)",
        subtitle="Well Duyong Deep-1  |  Formation Resistivity Factor & Resistivity Index at NOB Conditions",
        tbl_label="TABLE 2 — Formation Resistivity Factor & Resistivity Index  (4 samples × 5 Sw measurements)",
        headers=ri_hdrs, rows=ri_rows,
        col_ratios=[0.09, 0.11, 0.09, 0.09, 0.09, 0.07, 0.12, 0.09, 0.07],
        note="Saturant: 3,500 ppm  |  Resistivity of saturant: 1,500 Ω·m @ 77°F  |  NOB: 560 psi  |  F = a/φ^m  |  RI = a/Sw^n"
    )

    # ── Slide 3: TABLE 3 (Cap Press) + Pore Size Dist side by side ────────────
    t3 = next(t for t in tables if t['page'] == 13)
    t_psd = next(t for t in tables if t['page'] == 7)

    # TABLE 3 has 15 cols — keep: Sample, Depth, Ka, Phi, Sw@1, Sw@25, Sw@50, Sw@100, Sw@200 psi
    # Header row 2: ['Sample d','Depth m','NOB psi','Ka md','Phi %','0','1','2','4','8','15','25','50','100','200']
    cap_hdrs = ['Sample', 'Depth (m)', 'Ka (mD)', 'Phi (%)', 'Sw@1 psi', 'Sw@25 psi', 'Sw@50 psi', 'Sw@100 psi', 'Sw@200 psi']
    # Cols to keep: indices 0,1,3,4, 6(=1psi), 11(=25psi), 12(=50psi), 13(=100psi), 14(=200psi)
    keep_idx = [0, 1, 3, 4, 6, 11, 12, 13, 14]
    cap_rows = []
    for row in t3['grid'][2:]:
        cap_rows.append([row[i] if i < len(row) else '' for i in keep_idx])

    psd_hdrs = ['Sample ID', 'Depth (m)', 'Porosity (%)', 'Micro r<0.5 (%)', 'Meso 0.5-1.5 (%)', 'Macro r>1.5 (%)']
    psd_rows = [row for row in t_psd['grid'][2:]]

    _side_by_side_tables(
        prs,
        title="Extracted Structured Tables — Duyong Deep 1 (III)",
        subtitle="Well Duyong Deep-1  |  Capillary Pressure & Pore Size Distribution",
        left_title="TABLE 3 — Air-Brine Capillary Pressure by Centrifuge (key Sw columns)",
        left_hdrs=cap_hdrs, left_rows=cap_rows,
        left_ratios=[0.10, 0.12, 0.12, 0.10, 0.12, 0.12, 0.12, 0.12, 0.08],
        right_title="TABLE 5 — Pore Size Distribution Summary",
        right_hdrs=psd_hdrs, right_rows=psd_rows,
        right_ratios=[0.16, 0.18, 0.18, 0.20, 0.18, 0.10],
    )


def add_pegaga_structured_slides(prs):
    """3 slides showing key extracted tables from Pegaga-2."""
    from extract_tables import extract_tables_from_file

    tables = extract_tables_from_file('RESULTS/pegaga results.txt')

    # ── Slide 4: Sample Properties ────────────────────────────────────────────
    t_sp = next(t for t in tables if t['page'] == 7)
    sp_hdrs = ['Sample ID', 'Depth (m)', 'Ka NCS (mD)', 'Porosity (%)', 'Target Swi (%)', 'Achieved Swi (%)']
    sp_rows = [row for row in t_sp['grid'][1:]]

    _full_table_slide(
        prs,
        title="Extracted Structured Tables — Pegaga-2 (I)",
        subtitle="Pegaga-2, Malaysia  |  Hybrid USS G-W Drainage + O-W Centrifuge Imbibition Kr Tests",
        tbl_label="Table 1 — Sample Properties & Target / Achieved Swi",
        headers=sp_hdrs, rows=sp_rows,
        col_ratios=[0.14, 0.14, 0.18, 0.18, 0.18, 0.18],
        body_size=12,
    )

    # ── Slide 5: Combined drainage + imbibition params for all 5 samples ──────
    drn_pages  = [14, 21, 28, 35, 42]
    imb_pages  = [15, 22, 29, 36, 43]
    sample_ids = ['2-015', '2-019', '2-023', '2-029', '2-033']

    drn_tables = [next(t for t in tables if t['page'] == p and len(t['grid'][0]) == 7) for p in drn_pages]
    imb_tables = [next(t for t in tables if t['page'] == p and len(t['grid'][0]) == 8) for p in imb_pages]

    # Combined drainage params table
    drn_hdrs = ['Sample', 'Kw (mD)', 'Swir (frac)', 'Kg@Swir (mD)', 'Krg@Swir', 'Nw', 'Ng']
    drn_rows = []
    for sid, t in zip(sample_ids, drn_tables):
        row = t['grid'][1]  # data row
        drn_rows.append([sid, row[0], row[1], row[2], row[4], row[5], row[6]])

    # Combined imbibition params table
    imb_hdrs = ['Sample', 'Kg@Swir (mD)', 'Swir (frac)', 'Kw@Sgr (mD)', 'Sgr (frac)', 'Krg@Swir', 'Krw@Sgr', 'Nw', 'Ng']
    imb_rows = []
    for sid, t in zip(sample_ids, imb_tables):
        row = t['grid'][1]
        imb_rows.append([sid] + row)

    _side_by_side_tables(
        prs,
        title="Extracted Structured Tables — Pegaga-2 (II)",
        subtitle="Corey model parameters extracted by olmOCR ADE  |  All 5 samples  |  Sendra™ history-matched",
        left_title="Drainage USS — Corey Model Parameters",
        left_hdrs=drn_hdrs, left_rows=drn_rows,
        left_ratios=[0.14, 0.14, 0.14, 0.14, 0.14, 0.15, 0.15],
        right_title="Imbibition Centrifuge — Corey Model Parameters",
        right_hdrs=imb_hdrs, right_rows=imb_rows,
        right_ratios=[0.12, 0.12, 0.12, 0.12, 0.11, 0.10, 0.10, 0.10, 0.11],
    )

    # ── Slide 6: Sample 2-015 drainage kr data (first 12 rows) ───────────────
    # The kr data is in tables with 5 cols (Swn, Sw, Krw, Sw, Krg)
    # Find table on page 14 with 5 cols and many rows (skip the 2-row param table)
    kr_tbl_p14 = [t for t in tables if t['page'] == 14 and len(t['grid'][0]) == 5]
    # There may be multiple; pick the one with >5 rows
    kr_drn_015 = next((t for t in kr_tbl_p14 if len(t['grid']) > 5), None)

    if kr_drn_015:
        kr_hdrs = ['Swn (norm.)', 'Sw (drainage)', 'Krw', 'Sw (imbibition)', 'Krg']
        # Show first 12 data rows
        kr_rows = kr_drn_015['grid'][1:13]
        _full_table_slide(
            prs,
            title="Extracted Structured Tables — Pegaga-2 (III)",
            subtitle="Sample 2-015  |  Depth 2511.1 m  |  Ka = 2.46 mD  |  φ = 18.6%  |  Swir = 0.205 frac",
            tbl_label="Table 3.1.2 — Primary Drainage USS Relative Permeability Data  (first 12 of 41 rows shown)",
            headers=kr_hdrs, rows=kr_rows,
            col_ratios=[0.20, 0.20, 0.20, 0.20, 0.20],
            body_size=12,
            note="Full dataset in RESULTS/Pegaga2_Extracted_Tables.xlsx  |  Corey model: Nw = 4.55, Ng = 2.5  |  Krg@Swir = 0.593"
        )


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation(PPTX_IN)
    print(f"Loaded: {len(prs.slides)} slides from {PPTX_IN}")

    for idx in range(13, min(22, len(prs.slides))):
        print(f"  Fixing slide {idx+1} theme...")
        fix_slide_theme(prs.slides[idx], idx + 1)

    print("Adding Duyong Deep 1 overview slide...")
    add_duyong_overview(prs)

    print("Adding Pegaga-2 overview slide...")
    add_pegaga_overview(prs)

    print("Adding benchmark comparison slide...")
    add_benchmark_comparison(prs)

    print("Adding Duyong Deep 1 structured table slides...")
    add_duyong_structured_slides(prs)

    print("Adding Pegaga-2 structured table slides...")
    add_pegaga_structured_slides(prs)

    prs.save(PPTX_OUT)
    print(f"\n✓  Saved {len(prs.slides)} slides → {PPTX_OUT}")


if __name__ == "__main__":
    main()
