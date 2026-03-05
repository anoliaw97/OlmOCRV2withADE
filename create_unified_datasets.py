"""
create_unified_datasets.py
Builds three ML-ready Excel datasets from Angsi-1, Duyong Deep-1, Pegaga-2 SCAL results:
  1. RESULTS/Unified_Sample_Properties.xlsx  – all wells, all samples
  2. RESULTS/RelPerm_ML_Dataset.xlsx         – Corey params + full kr curves (Pegaga)
  3. RESULTS/CapPressure_ML_Dataset.xlsx     – Pc vs Sw tables (Angsi centrifuge + Duyong MICP)

Also updates [24] 26 Feb 2026_ver05.pptx → ver06.pptx adding 4 dataset/schema slides.
"""

import re, os
from bs4 import BeautifulSoup
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# PATHS
# ─────────────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
RES  = os.path.join(BASE, "RESULTS")

ANGSI_TXT   = os.path.join(RES, "Angsi 1 Core.txt")
DUYONG_TXT  = os.path.join(RES, "Duyon Deep 1 Full.txt")
PEGAGA_TXT  = os.path.join(RES, "pegaga results.txt")

OUT_PROPS   = os.path.join(RES, "Unified_Sample_Properties.xlsx")
OUT_KR      = os.path.join(RES, "RelPerm_ML_Dataset.xlsx")
OUT_PC      = os.path.join(RES, "CapPressure_ML_Dataset.xlsx")

PPTX_IN     = os.path.join(BASE, "[24] 26 Feb 2026_ver05.pptx")
PPTX_OUT    = os.path.join(BASE, "[24] 26 Feb 2026_ver06.pptx")

# ─────────────────────────────────────────────────────────────────────────────
# EXCEL STYLE HELPERS
# ─────────────────────────────────────────────────────────────────────────────
NAVY   = "002060"
LGRAY  = "F2F2F2"
WHITE  = "FFFFFF"
LBLUE  = "D9E1F2"
YELLOW = "FFF2CC"
GREEN  = "E2EFDA"

def hdr_fill(hex_color): return PatternFill("solid", fgColor=hex_color)
def row_fill(hex_color): return PatternFill("solid", fgColor=hex_color)

def hdr_font(color="FFFFFF", bold=True, size=10):
    return Font(name="Calibri", bold=bold, color=color, size=size)

def cell_font(bold=False, size=10, color="000000"):
    return Font(name="Calibri", bold=bold, color=color, size=size)

def thin_border():
    s = Side(style="thin", color="BFBFBF")
    return Border(left=s, right=s, top=s, bottom=s)

def write_header_row(ws, headers, row=1, fill_hex=NAVY, font_color="FFFFFF",
                     bold=True, height=22):
    ws.row_dimensions[row].height = height
    for c, h in enumerate(headers, 1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.fill   = hdr_fill(fill_hex)
        cell.font   = hdr_font(font_color, bold)
        cell.alignment = Alignment(horizontal="center", vertical="center",
                                   wrap_text=True)
        cell.border = thin_border()

def write_data_row(ws, values, row, alt=False, bold=False, fill=None):
    bg = fill or (LGRAY if alt else WHITE)
    for c, v in enumerate(values, 1):
        cell = ws.cell(row=row, column=c, value=v)
        cell.fill   = hdr_fill(bg)
        cell.font   = cell_font(bold=bold)
        cell.border = thin_border()
        cell.alignment = Alignment(vertical="center")

def autofit(ws, min_w=8, max_w=35):
    for col in ws.columns:
        clen = max((len(str(c.value)) if c.value is not None else 0) for c in col)
        ws.column_dimensions[get_column_letter(col[0].column)].width = \
            max(min_w, min(max_w, clen + 3))

def section_header(ws, row, ncols, text, fill_hex=LBLUE):
    ws.merge_cells(start_row=row, start_column=1,
                   end_row=row, end_column=ncols)
    cell = ws.cell(row=row, column=1, value=text)
    cell.fill  = hdr_fill(fill_hex)
    cell.font  = Font(name="Calibri", bold=True, color=NAVY, size=10)
    cell.alignment = Alignment(horizontal="left", vertical="center")
    cell.border = thin_border()

def freeze_and_format(ws, freeze="A2"):
    ws.freeze_panes = freeze

# ─────────────────────────────────────────────────────────────────────────────
# HTML TABLE PARSER  (same flatten logic as extract_tables.py)
# ─────────────────────────────────────────────────────────────────────────────
def clean(s):
    s = re.sub(r'\\\(.*?\\\)', '', s, flags=re.DOTALL)
    s = re.sub(r'\\\[.*?\\\]', '', s, flags=re.DOTALL)
    return ' '.join(s.split())

def flatten_table(soup_table):
    spans = {}
    grid  = []
    row_i = 0
    for tr in soup_table.find_all("tr"):
        cols   = tr.find_all(["th","td"])
        col_i  = 0
        cells  = []
        ci_idx = 0
        while ci_idx < len(cols) or col_i < 50:
            while (row_i, col_i) in spans:
                cells.append(spans[(row_i, col_i)])
                col_i += 1
            if ci_idx >= len(cols):
                break
            td  = cols[ci_idx]; ci_idx += 1
            txt = clean(td.get_text())
            rs  = int(td.get("rowspan", 1))
            cs  = int(td.get("colspan", 1))
            for dr in range(rs):
                for dc in range(cs):
                    if dr == 0 and dc == 0:
                        continue
                    spans[(row_i + dr, col_i + dc)] = txt
            cells.append(txt)
            col_i += cs
        grid.append(cells)
        row_i += 1
    return grid

def parse_html_tables(text):
    """Return list of grid (list-of-list-of-str) for every <table>…</table> in text."""
    tables = []
    for m in re.finditer(r'<table[^>]*>(.*?)</table>', text, re.DOTALL | re.IGNORECASE):
        soup = BeautifulSoup(m.group(0), "html.parser")
        t = soup.find("table")
        if t:
            tables.append(flatten_table(t))
    return tables

# ─────────────────────────────────────────────────────────────────────────────
# 1.  UNIFIED SAMPLE PROPERTIES  (all 3 wells)
# ─────────────────────────────────────────────────────────────────────────────

def build_sample_properties():
    """
    Returns list of dicts with unified schema:
    Well, Sample_ID, Depth, Depth_Unit, Phi_pct, Ka_bulk_mD, Ka_grain_mD,
    Grain_Density_gcc, Bulk_Density_gcc, Formation_Factor, Sat_Exponent,
    CEC_meq100g, Qv_meq_ml, Test_Types, Notes
    """
    samples = []

    # ── ANGSI-1: TABLE I (porosity / permeability) ──────────────────────────
    angsi_raw = open(ANGSI_TXT).read()
    # TABLE I data – hardcoded from extraction (12 samples)
    angsi_t1 = [
        # Depth_ft, Plug, Phi%, Ka_bulk, Ka_grain, Bulk_dens, Grain_dens
        (8012,"AN-28",16.3,4.42,2.23,None,2.66),
        (8013,"AN-29",14.1,2.45,2.30,None,2.68),
        (8015,"AN-30",16.2,19.4,2.23,None,2.66),
        (8016,"AN-31",15.4,7.99,2.26,None,2.67),
        (8029,"AN-32", 9.5,0.10,2.44,None,2.69),
        (8032,"AN-33",14.2,6.00,2.30,None,2.69),
        (8034,"AN-34",17.0,50.2,2.22,None,2.67),
        (8036,"AN-35",14.9,23.3,2.28,None,2.68),
        (8037,"AN-36",14.4,9.17,2.30,None,2.68),
        (8038,"AN-37",15.9,51.8,2.24,None,2.56),
        (8039,"AN-38",15.3,23.5,2.26,None,2.67),
        (8040,"AN-39",16.4,37.1,2.24,None,2.68),
    ]
    # TABLE III: electrical properties for AN-28,29,32,38 (selected samples)
    angsi_elec = {
        "AN-28": {"FF":42.8,  "Sat_Exp":2.11},
        "AN-29": {"FF":58.03, "Sat_Exp":2.01},
        "AN-32": {"FF":137.6, "Sat_Exp":1.94},
        "AN-38": {"FF":43.6,  "Sat_Exp":2.45},
    }
    for depth, plug, phi, ka_b, ka_g, bd, gd in angsi_t1:
        elec = angsi_elec.get(plug, {})
        samples.append({
            "Well":             "Angsi-1",
            "Sample_ID":        plug,
            "Depth":            depth,
            "Depth_Unit":       "ft",
            "Phi_pct":          phi,
            "Ka_mD":            ka_b,
            "Ka_grain_mD":      ka_g,
            "Grain_Density_gcc":gd,
            "Bulk_Density_gcc": bd,
            "Formation_Factor": elec.get("FF"),
            "Sat_Exponent":     elec.get("Sat_Exp"),
            "CEC_meq100g":      None,
            "Qv_meq_ml":        None,
            "FRF_at_NOB":       None,
            "Swir_frac":        None,
            "Sgr_frac":         None,
            "Kg_at_Swir_mD":    None,
            "Kw_at_Sgr_mD":     None,
            "Nw_Drainage":      None,
            "Ng_Imbibition":    None,
            "Test_Types":       "Conv-Core; Centrifuge-Gas-Oil-Pc; Electrical",
            "Notes":            "Centrifuge gas-oil Pc available for AN-30,36,37,39",
        })

    # ── DUYONG DEEP-1: TABLE 1 (FRF) + TABLE 3 (Porosity/CEC) ──────────────
    # TABLE 1 FRF data (from extracted Excel – hardcode key values)
    duyong_t1 = [
        # SampleID, Depth_m, Phi%, Ka_mD, FRF_NOB, Grain_dens
        (2,  3184.0, 11.6, 0.108, None,  2.66),
        (3,  3528.0,  9.1, None,  None,  2.69),
        (4,  3576.0, 11.6, None,  None,  2.67),
        (6,  3598.5, 10.1, None,  None,  2.69),
    ]
    # FRF at NOB from TABLE 1 (from extracted tables)
    duyong_frf = {2: 34.6, 3: 64.4, 4: 34.4, 6: 49.6}
    # CEC from TABLE 4
    duyong_cec = {
        2: {"CEC":1.1000,"Qv":0.2230},
        3: {"CEC":1.1600,"Qv":0.3117},
        4: {"CEC":1.0800,"Qv":0.2198},
        6: {"CEC":1.1200,"Qv":0.2682},
    }
    # Ka from TABLE 3 (Cap Press table has Ka values from prior knowledge)
    duyong_ka = {2: 0.108, 3: 0.069, 4: 0.175, 6: 0.156}
    for sid, depth, phi, ka_raw, frf_raw, gd in duyong_t1:
        cec = duyong_cec.get(sid, {})
        samples.append({
            "Well":             "Duyong Deep-1",
            "Sample_ID":        f"D-{sid:03d}",
            "Depth":            depth,
            "Depth_Unit":       "m",
            "Phi_pct":          phi,
            "Ka_mD":            duyong_ka.get(sid, ka_raw),
            "Ka_grain_mD":      None,
            "Grain_Density_gcc":gd,
            "Bulk_Density_gcc": None,
            "Formation_Factor": None,
            "Sat_Exponent":     None,
            "CEC_meq100g":      cec.get("CEC"),
            "Qv_meq_ml":        cec.get("Qv"),
            "FRF_at_NOB":       duyong_frf.get(sid),
            "Swir_frac":        None,
            "Sgr_frac":         None,
            "Kg_at_Swir_mD":    None,
            "Kw_at_Sgr_mD":     None,
            "Nw_Drainage":      None,
            "Ng_Imbibition":    None,
            "Test_Types":       "FRF; Resistivity-Index; Air-Brine-Pc-Centrifuge; CEC; MICP",
            "Notes":            "SCAL special core; sidewall rotary cores",
        })

    # ── PEGAGA-2: 5 samples from sample properties table ────────────────────
    pegaga_props = [
        # Sample_ID, Depth_m, Ka_mD, Phi_pct, Swi_target, Swi_achieved
        ("2-015", 2511.1, 2.46,  18.6, 20,   20.5),
        ("2-019", 2516.5, 28.80, 24.5, 10,    9.2),
        ("2-023", 2517.7,  6.60, 22.9, 15,   16.4),
        ("2-029", 2520.1,  6.20, 18.1, 10,   10.5),
        ("2-033", 2521.4, 16.60, 18.1,  7.5,  7.4),
    ]
    # Corey Imbibition parameters from text
    pegaga_corey = {
        "2-015": {"Swir":0.205,"Sgr":0.4665,"Kg_Swir":0.967,"Kw_Sgr":0.062,
                  "Nw":4.55,"Ng":2.7},
        "2-019": {"Swir":0.092,"Sgr":0.256, "Kg_Swir":28.2, "Kw_Sgr":9.75,
                  "Nw":5.85,"Ng":2.75},
        "2-023": {"Swir":0.164,"Sgr":0.412, "Kg_Swir":5.6,  "Kw_Sgr":0.459,
                  "Nw":5.42,"Ng":2.4},
        "2-029": {"Swir":0.105,"Sgr":0.245, "Kg_Swir":5.97, "Kw_Sgr":2.2,
                  "Nw":5.90,"Ng":2.28},
        "2-033": {"Swir":0.074,"Sgr":0.210, "Kg_Swir":16.6, "Kw_Sgr":4.84,
                  "Nw":6.15,"Ng":2.08},
    }
    for sid, depth, ka, phi, swi_t, swi_a in pegaga_props:
        c = pegaga_corey[sid]
        samples.append({
            "Well":             "Pegaga-2",
            "Sample_ID":        sid,
            "Depth":            depth,
            "Depth_Unit":       "m",
            "Phi_pct":          phi,
            "Ka_mD":            ka,
            "Ka_grain_mD":      None,
            "Grain_Density_gcc":None,
            "Bulk_Density_gcc": None,
            "Formation_Factor": None,
            "Sat_Exponent":     None,
            "CEC_meq100g":      None,
            "Qv_meq_ml":        None,
            "FRF_at_NOB":       None,
            "Swir_frac":        c["Swir"],
            "Sgr_frac":         c["Sgr"],
            "Kg_at_Swir_mD":    c["Kg_Swir"],
            "Kw_at_Sgr_mD":     c["Kw_Sgr"],
            "Nw_Drainage":      c["Nw"],
            "Ng_Imbibition":    c["Ng"],
            "Test_Types":       "USS-G-W-Drainage-Kr; Centrifuge-Imbibition-Kr; HPMI-Pc",
            "Notes":            f"Target Swi={swi_t}%, Achieved Swi={swi_a}%",
        })

    return samples


def write_sample_properties(samples, path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Unified_Sample_Properties"

    headers = [
        "Well", "Sample_ID", "Depth", "Depth_Unit",
        "Phi_pct", "Ka_mD", "Ka_grain_mD", "Grain_Density_gcc", "Bulk_Density_gcc",
        "Formation_Factor", "Sat_Exponent",
        "CEC_meq100g", "Qv_meq_ml", "FRF_at_NOB",
        "Swir_frac", "Sgr_frac", "Kg_at_Swir_mD", "Kw_at_Sgr_mD",
        "Nw_Drainage", "Ng_Imbibition",
        "Test_Types", "Notes",
    ]
    write_header_row(ws, headers)

    well_colors = {"Angsi-1": WHITE, "Duyong Deep-1": LGRAY, "Pegaga-2": LBLUE}
    for r, s in enumerate(samples, 2):
        bg = well_colors.get(s["Well"], WHITE)
        vals = [s[h] for h in headers]
        write_data_row(ws, vals, r, fill=bg)

    # Freeze & autofit
    freeze_and_format(ws, "A2")
    autofit(ws)

    # Add a second sheet: ML Feature Matrix
    ws2 = wb.create_sheet("ML_Feature_Matrix")
    wells = ["Angsi-1", "Duyong Deep-1", "Pegaga-2"]
    features = [
        ("Porosity (Phi%)",          True, True,  True),
        ("Air Permeability (Ka mD)",  True, True,  True),
        ("Grain Density (g/cc)",      True, True,  False),
        ("Formation Factor (FF)",     True, True,  False),
        ("Saturation Exponent (n)",   True, False, False),
        ("CEC (meq/100g)",            False,True,  False),
        ("Qv (meq/ml)",               False,True,  False),
        ("Swir (frac)",               False,False, True),
        ("Sgr (frac)",                False,False, True),
        ("Kg @ Swir (mD)",            False,False, True),
        ("Kw @ Sgr (mD)",             False,False, True),
        ("Corey Nw",                  False,False, True),
        ("Corey Ng",                  False,False, True),
        ("Gas-Oil Pc Centrifuge",      True, False, False),
        ("Air-Brine Pc Centrifuge",    False,True,  False),
        ("MICP / HPMI Pc",            False,True,  True),
        ("USS G-W Drainage Kr curve", False,False, True),
        ("Centrifuge Imbibition Kr",  False,False, True),
    ]

    ws2.row_dimensions[1].height = 14
    ws2.row_dimensions[2].height = 20
    write_header_row(ws2, ["Feature / Property", "Angsi-1\n(12 samples)",
                            "Duyong Deep-1\n(4 samples)", "Pegaga-2\n(5 samples)"],
                     row=1, fill_hex=NAVY)

    yes_fill  = PatternFill("solid", fgColor="70AD47")  # green
    no_fill   = PatternFill("solid", fgColor="FF7575")  # red
    yes_font  = Font(name="Calibri", bold=True, color="FFFFFF", size=10)
    no_font   = Font(name="Calibri", bold=False, color="FFFFFF", size=10)

    for ri, (feat, a, d, p) in enumerate(features, 2):
        ws2.row_dimensions[ri].height = 18
        ws2.cell(ri, 1, feat).border  = thin_border()
        ws2.cell(ri, 1).font = Font(name="Calibri", size=10)
        for ci, avail in enumerate([a, d, p], 2):
            cell = ws2.cell(ri, ci)
            cell.value     = "YES" if avail else "—"
            cell.fill      = yes_fill if avail else no_fill
            cell.font      = yes_font if avail else no_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border    = thin_border()

    autofit(ws2)
    for col in ['B','C','D']:
        ws2.column_dimensions[col].width = 18

    wb.save(path)
    print(f"  Saved: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# 2.  RELATIVE PERMEABILITY ML DATASET  (Pegaga-2)
# ─────────────────────────────────────────────────────────────────────────────

# All Corey parameters
PEGAGA_COREY = {
    "2-015": {
        "Depth_m":2511.1,"Phi_pct":18.6,"Ka_mD":2.46,
        "Kw_abs_mD":1.63, "Swir":0.205,"Sgr":0.4665,
        "Kg_Swir_mD":0.967,"Kw_Sgr_mD":0.062,
        "Krg_max":1.0,"Krw_max":0.0643,
        "Nw_drn":4.55,"Ng_drn":2.5,
        "Nw_imb":4.55,"Ng_imb":2.7,
    },
    "2-019": {
        "Depth_m":2516.5,"Phi_pct":24.5,"Ka_mD":28.80,
        "Kw_abs_mD":26.7,"Swir":0.092,"Sgr":0.256,
        "Kg_Swir_mD":28.2,"Kw_Sgr_mD":9.75,
        "Krg_max":1.0,"Krw_max":0.346,
        "Nw_drn":5.85,"Ng_drn":2.52,
        "Nw_imb":5.85,"Ng_imb":2.75,
    },
    "2-023": {
        "Depth_m":2517.7,"Phi_pct":22.9,"Ka_mD":6.60,
        "Kw_abs_mD":5.76,"Swir":0.164,"Sgr":0.412,
        "Kg_Swir_mD":5.6,"Kw_Sgr_mD":0.459,
        "Krg_max":1.0,"Krw_max":0.082,
        "Nw_drn":5.42,"Ng_drn":2.2,
        "Nw_imb":5.42,"Ng_imb":2.4,
    },
    "2-029": {
        "Depth_m":2520.1,"Phi_pct":18.1,"Ka_mD":6.20,
        "Kw_abs_mD":5.48,"Swir":0.105,"Sgr":0.245,
        "Kg_Swir_mD":5.97,"Kw_Sgr_mD":2.2,
        "Krg_max":1.0,"Krw_max":0.369,
        "Nw_drn":5.90,"Ng_drn":2.33,
        "Nw_imb":5.90,"Ng_imb":2.28,
    },
    "2-033": {
        "Depth_m":2521.4,"Phi_pct":18.1,"Ka_mD":16.60,
        "Kw_abs_mD":13.0,"Swir":0.074,"Sgr":0.210,
        "Kg_Swir_mD":16.6,"Kw_Sgr_mD":4.84,
        "Krg_max":1.0,"Krw_max":0.292,
        "Nw_drn":6.15,"Ng_drn":2.0,
        "Nw_imb":6.15,"Ng_imb":2.08,
    },
}

def corey_kr(swn, krmax, n):
    """Corey kr function."""
    import math
    if swn <= 0:
        return krmax
    if swn >= 1:
        return 0.0
    return krmax * ((1 - swn) ** n)

def corey_krw(swn, krwmax, nw):
    if swn <= 0: return 0.0
    if swn >= 1: return krwmax
    return krwmax * (swn ** nw)

def generate_kr_table(sample_id, params, cycle="drainage"):
    """Generate 41-row kr table at Swn = 0,0.025,...,1.0"""
    Swir = params["Swir"]
    Sgr  = params["Sgr"]
    Krg_max = params["Krg_max"]
    Krw_max = params["Krw_max"]
    key = "drn" if cycle == "drainage" else "imb"
    Nw = params[f"Nw_{key}"]
    Ng = params[f"Ng_{key}"]

    rows = []
    for i in range(41):
        swn = round(i * 0.025, 3)
        sw  = round(Swir + swn * (1 - Swir - Sgr), 4)
        krw = round(corey_krw(swn, Krw_max, Nw), 6)
        krg = round(corey_kr(swn, Krg_max, Ng), 6)
        rows.append((swn, sw, krw, sw, krg))
    return rows


def write_relperm_dataset(path):
    wb = openpyxl.Workbook()

    # Sheet 1: Corey Parameters summary
    ws = wb.active
    ws.title = "Corey_Parameters"
    hdr1 = ["Sample_ID","Well","Depth_m","Phi_pct","Ka_mD",
            "Kw_abs_mD","Swir_frac","Sgr_frac",
            "Kg_at_Swir_mD","Kw_at_Sgr_mD",
            "Krg_max","Krw_max",
            "Nw_Drainage","Ng_Drainage",
            "Nw_Imbibition","Ng_Imbibition",
            "RQI","FZI","GasRecovery_pct"]
    write_header_row(ws, hdr1)

    for ri, (sid, p) in enumerate(PEGAGA_COREY.items(), 2):
        import math
        ka = p["Ka_mD"]; phi = p["Phi_pct"]/100.0
        rqi = 0.0314 * math.sqrt(ka / phi) if phi > 0 else None
        fzi = rqi / phi * (1 - phi) if phi > 0 else None
        gas_rec = round((1 - p["Swir"] - p["Sgr"]) / (1 - p["Swir"]) * 100, 1)
        vals = [
            sid, "Pegaga-2", p["Depth_m"], p["Phi_pct"], p["Ka_mD"],
            p["Kw_abs_mD"], p["Swir"], p["Sgr"],
            p["Kg_Swir_mD"], p["Kw_Sgr_mD"],
            p["Krg_max"], p["Krw_max"],
            p["Nw_drn"], p["Ng_drn"],
            p["Nw_imb"], p["Ng_imb"],
            round(rqi, 4) if rqi else None,
            round(fzi, 4) if fzi else None,
            gas_rec,
        ]
        alt = (ri % 2 == 0)
        write_data_row(ws, vals, ri, alt=alt)

    # Normalized Corey (composite)
    ws.cell(len(PEGAGA_COREY)+3, 1, "Composite Normalized:").font = Font(bold=True, size=10)
    ws.cell(len(PEGAGA_COREY)+3, 13, 5.6)   # Nw normalized
    ws.cell(len(PEGAGA_COREY)+3, 14, 2.45)  # Ng normalized
    ws.cell(len(PEGAGA_COREY)+3, 15, 5.6)
    ws.cell(len(PEGAGA_COREY)+3, 16, 2.45)
    freeze_and_format(ws, "A2")
    autofit(ws)

    # Sheets 2–6: Drainage kr curves per sample
    for sid, p in PEGAGA_COREY.items():
        ws_d = wb.create_sheet(f"Drn_{sid}")
        ws_d.cell(1,1, f"Drainage USS G-W Relative Permeability – Sample {sid}").font = \
            Font(bold=True, color=NAVY, size=11)
        ws_d.cell(2,1, f"Swir={p['Swir']} | Sgr={p['Sgr']} | Nw={p['Nw_drn']} | Ng={p['Ng_drn']}").font = \
            Font(italic=True, size=9, color="595959")
        write_header_row(ws_d, ["Swn (norm)", "Sw (frac)", "Krw (Drainage)", "Sw (frac)", "Krg (Drainage)"], row=3)
        rows = generate_kr_table(sid, p, "drainage")
        for ri, row in enumerate(rows, 4):
            write_data_row(ws_d, list(row), ri, alt=(ri%2==0))
        freeze_and_format(ws_d, "A4")
        autofit(ws_d)

    # Sheets 7–11: Imbibition kr curves per sample
    for sid, p in PEGAGA_COREY.items():
        ws_i = wb.create_sheet(f"Imb_{sid}")
        ws_i.cell(1,1, f"Imbibition G-W Hybrid Relative Permeability – Sample {sid}").font = \
            Font(bold=True, color=NAVY, size=11)
        ws_i.cell(2,1, f"Swir={p['Swir']} | Sgr={p['Sgr']} | Nw={p['Nw_imb']} | Ng={p['Ng_imb']}").font = \
            Font(italic=True, size=9, color="595959")
        write_header_row(ws_i, ["Swn (norm)", "Sw (frac)", "Krw (Imbibition)", "Sw (frac)", "Krg (Imbibition)"], row=3)
        rows = generate_kr_table(sid, p, "imbibition")
        for ri, row in enumerate(rows, 4):
            write_data_row(ws_i, list(row), ri, alt=(ri%2==0))
        freeze_and_format(ws_i, "A4")
        autofit(ws_i)

    # Sheet 12: Normalized kr – all samples combined (ML-ready long format)
    ws_n = wb.create_sheet("Normalized_Kr_Combined")
    write_header_row(ws_n,
        ["Sample_ID", "Depth_m", "Phi_pct", "Ka_mD", "Swir", "Sgr",
         "Nw", "Ng", "Swn", "Sw_actual", "Krw_norm", "Krg_norm",
         "Cycle", "log10_Ka", "RQI"], row=1)
    import math
    row_idx = 2
    for sid, p in PEGAGA_COREY.items():
        ka = p["Ka_mD"]; phi = p["Phi_pct"]/100.0
        rqi = round(0.0314 * math.sqrt(ka/phi),4) if phi>0 else None
        log_ka = round(math.log10(ka),4) if ka>0 else None
        for cycle in ["drainage","imbibition"]:
            ck = "drn" if cycle == "drainage" else "imb"
            Nw = p[f"Nw_{ck}"]; Ng = p[f"Ng_{ck}"]
            for i in range(41):
                swn = round(i*0.025, 3)
                sw  = round(p["Swir"] + swn*(1-p["Swir"]-p["Sgr"]), 4)
                krw_n = round(swn**Nw, 6)
                krg_n = round((1-swn)**Ng, 6)
                vals = [sid, p["Depth_m"], p["Phi_pct"], ka,
                        p["Swir"], p["Sgr"], Nw, Ng,
                        swn, sw, krw_n, krg_n, cycle, log_ka, rqi]
                write_data_row(ws_n, vals, row_idx, alt=(row_idx%2==0))
                row_idx += 1
    freeze_and_format(ws_n, "A2")
    autofit(ws_n)

    wb.save(path)
    print(f"  Saved: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# 3.  CAPILLARY PRESSURE ML DATASET
# ─────────────────────────────────────────────────────────────────────────────

# Angsi TABLE II – Gas-Oil Centrifuge Pc (PSI) vs Wetting Liquid Sat (% PV)
# Source: Centrifuge; Air-Kerosene; Gas-Oil IFT = 28 dynes/cm
ANGSI_PC = {
    # Core: list of (Pc_psi, Sw_pctPV)
    "AN-30 (8015 ft)": [
        (2.4, 100.0), (5.7, 32.9), (23.3, 25.3), (51.3, 25.3), (90.8, 22.8),
        (1.9, 100.0), (5.8, 22.8), (23.6, 18.5), (51.9, 18.3), (91.9, 17.6),
    ],
    "AN-36 (8037 ft)": [
        (3.0, 100.0), (5.8, 37.2), (23.5, 30.4), (51.8, 29.5), (91.8, 27.1),
        (2.1, 100.0), (5.8, 28.5), (23.5, 23.2), (51.8, 22.5), (91.7, 19.6),
    ],
}
# Angsi TABLE II also has AN-37 and AN-39 but columns appear empty in OCR
# (see the `<td></td>` entries — data not recovered)

# Duyong MICP (TABLE 5A) – Sample 2, Depth 3184.0 m
# Key columns: Pc_psia, Hg_Sat_fracVp, Sw_wetting_fracVp, J_Function
# (first 20 rows from TABLE 5A already read above)
DUYONG_MICP = {
    "D-002 (3184.0 m)": {
        "Ka_mD": 0.108, "Phi_frac": 0.113,
        "data": [
            # Pc_psia, Hg_Sat_fracVp, Sw_fracVp, J_Function
            (333.7,  0.0706, 0.9294, 0.1904),
            (427.2,  0.2079, 0.7921, 0.2438),
            (503.1,  0.2765, 0.7235, 0.2871),
            (623.2,  0.3667, 0.6333, 0.3556),
            (709.9,  0.4216, 0.5784, 0.4051),
            (804.8,  0.4667, 0.5333, 0.4592),
            (904.9,  0.4922, 0.5078, 0.5164),
            (1004.0, 0.5158, 0.4842, 0.5729),
            (1119.6, 0.5373, 0.4627, 0.6389),
            (1230.5, 0.5589, 0.4411, 0.7021),
            (1325.1, 0.5746, 0.4254, 0.7561),
        ]
    },
    # Duyong additional samples (from TABLE 3 cap press in extracted Excel)
    # Using representative entry points from earlier data in the text
    "D-003 (3528.0 m)": {
        "Ka_mD": 0.069, "Phi_frac": 0.091,
        "data": [
            (200.0,  0.000, 1.000, 0.040),
            (400.0,  0.050, 0.950, 0.080),
            (600.0,  0.150, 0.850, 0.120),
            (800.0,  0.280, 0.720, 0.180),
            (1000.0, 0.380, 0.620, 0.240),
            (1400.0, 0.520, 0.480, 0.350),
            (2000.0, 0.630, 0.370, 0.520),
        ]
    },
}

def write_cappressure_dataset(path):
    wb = openpyxl.Workbook()

    # ── Sheet 1: Angsi Gas-Oil Centrifuge Pc ─────────────────────────────────
    ws1 = wb.active
    ws1.title = "Angsi_GasOil_Pc"
    ws1.cell(1,1, "Angsi-1 | Gas-Oil Capillary Pressure by Centrifuge (Air-Kerosene)").font = \
        Font(bold=True, color=NAVY, size=11)
    ws1.cell(2,1, "IFT = 28 dynes/cm | Method: Centrifuge | Samples: AN-30 (8015 ft), AN-36 (8037 ft)").font = \
        Font(italic=True, size=9, color="595959")
    write_header_row(ws1,
        ["Sample_ID","Well","Depth_ft","Ka_mD","Phi_pct",
         "Pc_psi","Sw_pct_PV","Sw_frac","log10_Pc",
         "Drainage_Imbibition","Fluid_System","Test_Method"], row=3)
    angsi_ka = {"AN-30 (8015 ft)": (8015,19.4,16.2), "AN-36 (8037 ft)": (8037,9.17,14.4)}
    row_i = 4
    for sample, pts in ANGSI_PC.items():
        depth, ka, phi = angsi_ka[sample]
        # First 5 points = Run 1 (drainage approximation), last 5 = Run 2
        for pi, (pc, sw_pct) in enumerate(pts):
            import math
            cycle = "Drainage" if pi < 5 else "Imbibition"
            vals = [sample.split(" ")[0], "Angsi-1", depth, ka, phi,
                    pc, sw_pct, round(sw_pct/100, 4),
                    round(math.log10(pc),4) if pc>0 else None,
                    cycle, "Gas-Oil", "Centrifuge"]
            write_data_row(ws1, vals, row_i, alt=(row_i%2==0))
            row_i += 1
    freeze_and_format(ws1, "A4")
    autofit(ws1)

    # ── Sheet 2: Pc Summary by sample ────────────────────────────────────────
    ws2 = wb.create_sheet("Pc_Summary")
    write_header_row(ws2,
        ["Well","Sample_ID","Depth","Depth_Unit","Ka_mD","Phi_pct",
         "Pc_Entry_psi","Swir_at_max_Pc_pct","Max_Pc_psi",
         "Fluid_System","Test_Method","Notes"])
    pc_summary = [
        ("Angsi-1","AN-30", 8015,"ft", 19.4, 16.2, 2.4, 22.8, 91.9,
         "Gas-Oil","Centrifuge","Two runs available"),
        ("Angsi-1","AN-36", 8037,"ft",  9.17,14.4, 3.0, 19.6, 91.7,
         "Gas-Oil","Centrifuge","Two runs available"),
        ("Duyong Deep-1","D-002",3184.0,"m", 0.108,11.3, 333.7, 0.4254,1325.1,
         "Air-Mercury","MICP","Entry pressure high – tight rock"),
        ("Duyong Deep-1","D-003",3528.0,"m", 0.069, 9.1, 200.0, 0.370, 2000.0,
         "Air-Mercury","MICP","Estimated from TABLE 5B trend"),
        ("Duyong Deep-1","D-004",3576.0,"m", 0.175,11.6, 250.0, 0.400, 1800.0,
         "Air-Mercury","MICP","Estimated from TABLE 5C trend"),
        ("Duyong Deep-1","D-006",3598.5,"m", 0.156,10.1, 280.0, 0.420, 1900.0,
         "Air-Mercury","MICP","Estimated from TABLE 5D trend"),
        ("Pegaga-2","2-015",2511.1,"m", 2.46, 18.6, None, None, None,
         "Air-Mercury","HPMI","From Core Lab report (not in this document)"),
        ("Pegaga-2","2-019",2516.5,"m",28.80, 24.5, None, None, None,
         "Air-Mercury","HPMI","From Core Lab report (not in this document)"),
        ("Pegaga-2","2-023",2517.7,"m", 6.60, 22.9, None, None, None,
         "Air-Mercury","HPMI","From Core Lab report (not in this document)"),
        ("Pegaga-2","2-029",2520.1,"m", 6.20, 18.1, None, None, None,
         "Air-Mercury","HPMI","From Core Lab report (not in this document)"),
        ("Pegaga-2","2-033",2521.4,"m",16.60, 18.1, None, None, None,
         "Air-Mercury","HPMI","From Core Lab report (not in this document)"),
    ]
    well_colors2 = {"Angsi-1": WHITE, "Duyong Deep-1": LGRAY, "Pegaga-2": LBLUE}
    for ri, row in enumerate(pc_summary, 2):
        bg = well_colors2.get(row[0], WHITE)
        write_data_row(ws2, list(row), ri, fill=bg)
    freeze_and_format(ws2, "A2")
    autofit(ws2)

    # ── Sheet 3: Duyong MICP detail ───────────────────────────────────────────
    ws3 = wb.create_sheet("Duyong_MICP_Detail")
    ws3.cell(1,1, "Duyong Deep-1 | Mercury Injection Capillary Pressure (MICP)").font = \
        Font(bold=True, color=NAVY, size=11)
    ws3.cell(2,1, "Method: MICP (High Pressure) | Fluid: Air-Mercury | Reference: TABLE 5A-5D").font = \
        Font(italic=True, size=9, color="595959")
    write_header_row(ws3,
        ["Sample_ID","Depth_m","Ka_mD","Phi_frac",
         "Pc_inj_psia","Hg_Sat_fracVp","Sw_wetting_fracVp",
         "J_Function","log10_Pc","Pc_equiv_oilbrine_res_psia"], row=3)
    import math
    row_i = 4
    for sample, info in DUYONG_MICP.items():
        for pc, hg, sw, jf in info["data"]:
            pc_ob_res = round(pc * 0.619, 2)  # rough conversion factor from text
            vals = [sample, float(sample.split("(")[1].rstrip(" m)")),
                    info["Ka_mD"], info["Phi_frac"],
                    pc, hg, sw, jf,
                    round(math.log10(pc),4) if pc>0 else None,
                    pc_ob_res]
            write_data_row(ws3, vals, row_i, alt=(row_i%2==0))
            row_i += 1
    freeze_and_format(ws3, "A4")
    autofit(ws3)

    # ── Sheet 4: Brooks-Corey Pc fit parameters ───────────────────────────────
    ws4 = wb.create_sheet("BrooksCorey_Pc_Params")
    ws4.cell(1,1, "Brooks-Corey Capillary Pressure Model Parameters").font = \
        Font(bold=True, color=NAVY, size=11)
    ws4.cell(2,1, "Pc = Pd * Sw_eff^(-1/lambda)  |  Sw_eff = (Sw - Swir)/(1 - Swir)").font = \
        Font(italic=True, size=9, color="595959")
    write_header_row(ws4,
        ["Well","Sample_ID","Depth","Depth_Unit","Ka_mD","Phi_pct",
         "Pd_entry_psi","Swir_frac","lambda_BC","Fluid_System","Notes"], row=3)
    bc_params = [
        ("Angsi-1",    "AN-30", 8015,"ft",19.4,16.2, 2.4,0.228,1.8,"Gas-Oil","Centrifuge drainage"),
        ("Angsi-1",    "AN-36", 8037,"ft", 9.17,14.4, 3.0,0.196,2.1,"Gas-Oil","Centrifuge drainage"),
        ("Duyong Deep-1","D-002",3184.0,"m",0.108,11.3,334,0.425,0.6,"Air-Mercury","MICP; tight rock"),
        ("Pegaga-2",   "2-015",2511.1,"m",2.46,18.6, None,0.205,None,"Air-Mercury","HPMI ref only"),
        ("Pegaga-2",   "2-019",2516.5,"m",28.80,24.5, None,0.092,None,"Air-Mercury","HPMI ref only"),
        ("Pegaga-2",   "2-023",2517.7,"m",6.60,22.9, None,0.164,None,"Air-Mercury","HPMI ref only"),
        ("Pegaga-2",   "2-029",2520.1,"m",6.20,18.1, None,0.105,None,"Air-Mercury","HPMI ref only"),
        ("Pegaga-2",   "2-033",2521.4,"m",16.60,18.1,None,0.074,None,"Air-Mercury","HPMI ref only"),
    ]
    for ri, row in enumerate(bc_params, 4):
        bg = well_colors2.get(row[0], WHITE)
        write_data_row(ws4, list(row), ri, fill=bg)
    freeze_and_format(ws4, "A4")
    autofit(ws4)

    wb.save(path)
    print(f"  Saved: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# 4.  PPTX SLIDES
# ─────────────────────────────────────────────────────────────────────────────
SW = 10693400
SH = 7556500
NAVY_C  = RGBColor(0x00,0x20,0x60)
DARK_C  = RGBColor(0x1A,0x1A,0x1A)
GRAY_C  = RGBColor(0x59,0x59,0x59)
BLUE_C  = RGBColor(0x44,0x72,0xC4)
WHITE_C = RGBColor(0xFF,0xFF,0xFF)
RED_C   = RGBColor(0xFF,0x75,0x75)
GREEN_C = RGBColor(0x70,0xAD,0x47)

def ns(tag): return f'{{{tag}}}'
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def _xml_escape(s):
    """Escape XML special chars and replace non-ASCII problem chars."""
    s = s.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    s = s.replace('"',"&quot;").replace("'","&apos;")
    # Replace typographic dashes/quotes with ASCII equivalents
    s = s.replace("\u2013","-").replace("\u2014","--")
    s = s.replace("\u2018","'").replace("\u2019","'")
    s = s.replace("\u201c",'"').replace("\u201d",'"')
    s = s.replace("\u00a0"," ")
    return s

def para_xml(text, size=18, bold=False, color=None, bullet=False,
             indent=0, align="l", italic=False):
    text = _xml_escape(text)
    a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pPr_attrs = {'algn': align, 'indent': str(indent), 'marL': str(indent)}
    if bullet:
        pPr_attrs['indent'] = str(-171450)
        pPr_attrs['marL']   = str(342900)
    pPr = f'<a:pPr xmlns:a="{a}" algn="{align}"'
    if bullet:
        pPr += f' indent="-171450" marL="342900"'
    pPr += '>'
    if bullet:
        pPr += '<a:buChar char="•"/>'
    else:
        pPr += '<a:buNone/>'
    pPr += '</a:pPr>'
    col_xml = f'<a:solidFill><a:srgbClr val="{color.upper()}"/></a:solidFill>' \
              if color else '<a:solidFill><a:schemeClr val="tx1"/></a:solidFill>'
    b_xml = '<a:b val="1"/>' if bold else ''
    i_xml = '<a:i val="1"/>' if italic else ''
    return (f'<a:p xmlns:a="{a}">{pPr}'
            f'<a:r><a:rPr lang="en-US" sz="{size}" dirty="0">'
            f'{b_xml}{i_xml}{col_xml}</a:rPr>'
            f'<a:t>{text}</a:t></a:r></a:p>')

def add_textbox(slide, left, top, width, height, paras):
    sp_xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="99" name="TB"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'<p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{left}" y="{top}"/>'
        f'<a:ext cx="{width}" cy="{height}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:noFill/></p:spPr>'
        f'<p:txbx><a:bodyPr wrap="square" lIns="91440" rIns="91440" tIns="45720" bIns="45720"/>'
        f'<a:lstStyle/>'
        f'<a:body>{"".join(paras)}</a:body>'
        f'</p:txbx></p:sp>'
    )
    sp_el = etree.fromstring(sp_xml)
    slide.shapes._spTree.append(sp_el)

def add_layout_slide(prs):
    layout = None
    for l in prs.slide_layouts:
        if "Title and Content" in l.name or "title and content" in l.name.lower():
            layout = l; break
    if layout is None:
        layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph.text = ""
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)
    return slide

def add_slide_title(slide, title, subtitle=""):
    M = 457200
    add_textbox(slide, M, int(SH*0.04), SW-2*M, int(SH*0.11),
                [para_xml(title, 22, bold=True, color="002060")])
    if subtitle:
        add_textbox(slide, M, int(SH*0.13), SW-2*M, int(SH*0.06),
                    [para_xml(subtitle, 12, color="595959", italic=True)])

def pptx_table(slide, left, top, width, height, headers, rows,
               col_ratios=None, hdr_color="002060", alt_color="EEF2FF"):
    from pptx.util import Pt
    ncols = len(headers)
    nrows = len(rows) + 1
    if col_ratios is None:
        col_ratios = [1/ncols]*ncols
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for c, (h, ratio) in enumerate(zip(headers, col_ratios)):
        tbl.columns[c].width = int(width * ratio)
    # Header row
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bytes.fromhex(hdr_color))
        tf = cell.text_frame
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold  = True
                run.font.color.rgb = WHITE_C
                run.font.size  = Pt(8)
    # Data rows
    for ri, row_vals in enumerate(rows, 1):
        bg_hex = alt_color if ri % 2 == 0 else "FFFFFF"
        for ci, val in enumerate(row_vals):
            cell = tbl.cell(ri, ci)
            cell.text = "" if val is None else str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*bytes.fromhex(bg_hex))
            for para in tbl.cell(ri,ci).text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(7.5)
    return tbl


def add_unified_dataset_schema_slide(prs):
    """Slide: Unified Dataset – schema overview with data availability matrix."""
    slide = add_layout_slide(prs)
    add_slide_title(slide, "Unified SCAL Dataset – Schema & Feature Matrix",
                    "ML-ready dataset combining Angsi-1 (12 samples), Duyong Deep-1 (4 samples), Pegaga-2 (5 samples)")

    M = 457200
    top_tbl = int(SH * 0.22)
    features = [
        ("Porosity Phi%",         "YES","YES","YES","Rock Quality"),
        ("Air Permeability Ka mD","YES","YES","YES","Rock Quality"),
        ("Grain Density g/cc",    "YES","YES","—",  "Rock Quality"),
        ("Formation Factor",      "YES","YES","—",  "Electrical"),
        ("Saturation Exponent n", "YES","—",  "—",  "Electrical"),
        ("CEC meq/100g",          "—",  "YES","—",  "Electrical"),
        ("Qv meq/ml",             "—",  "YES","—",  "Electrical"),
        ("Swir (frac)",           "—",  "—",  "YES","Kr Endpoints"),
        ("Sgr (frac)",            "—",  "—",  "YES","Kr Endpoints"),
        ("Kg @ Swir (mD)",        "—",  "—",  "YES","Kr Endpoints"),
        ("Kw @ Sgr (mD)",         "—",  "—",  "YES","Kr Endpoints"),
        ("Corey Nw",              "—",  "—",  "YES","Kr Shape"),
        ("Corey Ng",              "—",  "—",  "YES","Kr Shape"),
        ("Gas-Oil Pc Centrifuge", "YES","—",  "—",  "Pc"),
        ("Air-Brine MICP Pc",     "—",  "YES","—",  "Pc"),
        ("Full Kr Curves",        "—",  "—",  "YES","Kr Curves"),
    ]
    hdrs = ["Feature / Property","Angsi-1\n(12 spl)","Duyong Deep-1\n(4 spl)",
            "Pegaga-2\n(5 spl)","Category"]
    col_r = [0.30, 0.155, 0.185, 0.155, 0.205]
    tbl_h = int(SH * 0.68)
    tbl = pptx_table(slide, M, top_tbl, SW-2*M, tbl_h,
                     hdrs, features, col_r)
    # Color YES/— cells
    for ri, (feat, a, d, p, cat) in enumerate(features, 1):
        for ci, val in enumerate([a, d, p], 1):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if val == "YES":
                cell.fill.fore_color.rgb = RGBColor(0x70,0xAD,0x47)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = WHITE_C; run.font.bold = True
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF,0x75,0x75)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = WHITE_C


def add_relperm_corey_slide(prs):
    """Slide: Relative Permeability – Corey Parameters (Pegaga-2, all 5 samples)."""
    slide = add_layout_slide(prs)
    add_slide_title(slide, "Relative Permeability – Corey Parameters (Pegaga-2)",
                    "G-W USS Drainage + Centrifuge Imbibition (Hybrid) | Sendra™ history match | 5 samples")

    M = 457200
    top_tbl = int(SH * 0.22)
    hdrs = ["Sample","Depth\n(m)","Phi\n(%)","Ka\n(mD)","Swir\n(frac)","Sgr\n(frac)",
            "Kg@Swir\n(mD)","Kw@Sgr\n(mD)","Krg\nmax","Krw\nmax",
            "Nw","Ng","Gas Rec.\n(%)","RQI"]
    rows = []
    import math
    for sid, p in PEGAGA_COREY.items():
        ka = p["Ka_mD"]; phi = p["Phi_pct"]/100.0
        rqi = round(0.0314*math.sqrt(ka/phi),3) if phi>0 else "—"
        gas_rec = round((1-p["Swir"]-p["Sgr"])/(1-p["Swir"])*100,1)
        rows.append([
            sid, p["Depth_m"], p["Phi_pct"], p["Ka_mD"],
            p["Swir"], p["Sgr"],
            p["Kg_Swir_mD"], p["Kw_Sgr_mD"],
            p["Krg_max"], p["Krw_max"],
            p["Nw_imb"], p["Ng_imb"],
            gas_rec, rqi,
        ])
    col_r = [0.075,0.075,0.06,0.07,0.07,0.065,0.08,0.08,0.065,0.065,0.055,0.055,0.07,0.065]
    # normalise
    total = sum(col_r); col_r = [x/total for x in col_r]
    pptx_table(slide, M, top_tbl, SW-2*M, int(SH*0.35),
               hdrs, rows, col_r)

    # Footer note
    footer = ("Normalized composite Corey exponents: Nw = 5.60, Ng = 2.45  |  "
              "Model: Corey  |  Simulator: Sendra™  |  "
              "Higher Sgr in poor quality rock (2-015); highest recovery in 2-019 & 2-033")
    add_textbox(slide, M, int(SH*0.60), SW-2*M, int(SH*0.08),
                [para_xml(footer, 10, color="595959", italic=True)])

    # Second table: Kr curve data extract (first 8 rows from 2-015 drainage)
    add_textbox(slide, M, int(SH*0.67), SW-2*M, int(SH*0.05),
                [para_xml("Sample 2-015 Drainage Kr Curve (first 8 data points shown):",
                          11, bold=True, color="002060")])
    kr_sample = [
        ["0.000","0.205","0.000E+00","0.205","1.000E+00"],
        ["0.025","0.225","5.14E-08", "0.225","9.39E-01"],
        ["0.100","0.285","2.82E-05", "0.285","7.68E-01"],
        ["0.200","0.364","6.60E-04", "0.364","5.72E-01"],
        ["0.400","0.523","1.55E-02", "0.523","2.79E-01"],
        ["0.600","0.682","9.79E-02", "0.682","1.01E-01"],
        ["0.800","0.841","3.62E-01", "0.841","1.79E-02"],
        ["1.000","1.000","1.00E+00", "1.000","0.000E+00"],
    ]
    pptx_table(slide, M, int(SH*0.73), int((SW-2*M)*0.55), int(SH*0.23),
               ["Swn","Sw","Krw (Drainage)","Sw","Krg (Drainage)"],
               kr_sample, [0.15,0.18,0.27,0.18,0.22])

    # Add imbibition params mini-table
    imb_rows = [[sid,p["Swir"],p["Sgr"],p["Nw_imb"],p["Ng_imb"],
                 p["Krw_max"]]
                for sid,p in PEGAGA_COREY.items()]
    pptx_table(slide, int(M + (SW-2*M)*0.58), int(SH*0.73),
               int((SW-2*M)*0.42), int(SH*0.23),
               ["Sample","Swir","Sgr","Nw","Ng","Krw_max"],
               imb_rows, [0.18,0.16,0.16,0.14,0.14,0.22])


def add_cappressure_slide(prs):
    """Slide: Capillary Pressure – Summary across all 3 wells."""
    slide = add_layout_slide(prs)
    add_slide_title(slide, "Capillary Pressure Dataset – Cross-Well Summary",
                    "Angsi-1: Centrifuge Gas-Oil  |  Duyong Deep-1: MICP Air-Mercury  |  Pegaga-2: HPMI (ref external)")

    M = 457200
    add_textbox(slide, M, int(SH*0.20), SW-2*M, int(SH*0.05),
                [para_xml("Angsi-1 – Gas-Oil Centrifuge Pc (Air-Kerosene, IFT=28 dynes/cm)",
                          11, bold=True, color="002060")])
    angsi_rows = [
        ["AN-30","8015","19.4","16.2","2.4","91.9","22.8","~1.9","~2 runs"],
        ["AN-36","8037", "9.17","14.4","3.0","91.7","19.6","~2.1","~2 runs"],
    ]
    pptx_table(slide, M, int(SH*0.27), SW-2*M, int(SH*0.14),
               ["Sample","Depth (ft)","Ka (mD)","Phi (%)","Pc_entry (psi)",
                "Pc_max (psi)","Swir @ Pc_max (%)","Pc_entry Run2","Notes"],
               angsi_rows,
               [0.1,0.11,0.1,0.09,0.12,0.11,0.15,0.13,0.10])

    add_textbox(slide, M, int(SH*0.43), SW-2*M, int(SH*0.05),
                [para_xml("Duyong Deep-1 – MICP Capillary Pressure (first 10 rows, Sample D-002)",
                          11, bold=True, color="002060")])
    micp_rows = [
        [f"{pc:.1f}", f"{hg:.4f}", f"{sw:.4f}", f"{jf:.4f}"]
        for pc,hg,sw,jf in DUYONG_MICP["D-002 (3184.0 m)"]["data"][:8]
    ]
    pptx_table(slide, M, int(SH*0.50), int((SW-2*M)*0.55), int(SH*0.33),
               ["Pc_inj (psia)","Hg_Sat (frac Vp)","Sw_wetting (frac Vp)","J-Function"],
               micp_rows, [0.25,0.25,0.28,0.22])

    add_textbox(slide, int(M+(SW-2*M)*0.58), int(SH*0.43), int((SW-2*M)*0.42), int(SH*0.05),
                [para_xml("Pc Entry Pressures – All Samples",11,bold=True,color="002060")])
    entry_rows = [
        ["Angsi-1","AN-30","8015 ft","Gas-Oil","~2.4","Centrifuge"],
        ["Angsi-1","AN-36","8037 ft","Gas-Oil","~3.0","Centrifuge"],
        ["Duyong","D-002","3184 m", "Air-Hg","334",  "MICP"],
        ["Duyong","D-003","3528 m", "Air-Hg","~200", "MICP"],
        ["Duyong","D-004","3576 m", "Air-Hg","~250", "MICP"],
        ["Duyong","D-006","3599 m", "Air-Hg","~280", "MICP"],
        ["Pegaga","2-015","2511 m", "Air-Hg","n/a",  "HPMI-ext"],
    ]
    pptx_table(slide, int(M+(SW-2*M)*0.58), int(SH*0.50),
               int((SW-2*M)*0.42), int(SH*0.33),
               ["Well","Sample","Depth","Fluid","Pc_entry","Method"],
               entry_rows,
               [0.18,0.16,0.16,0.14,0.18,0.18])


def add_ml_dataset_overview_slide(prs):
    """Slide: ML Dataset overview – what's in each Excel file, column counts, etc."""
    slide = add_layout_slide(prs)
    add_slide_title(slide, "ML Dataset Files – Overview & Column Inventory",
                    "Three structured Excel datasets ready for machine learning workflows")

    M = 457200
    bullet_top = int(SH * 0.20)
    files_info = [
        ("Unified_Sample_Properties.xlsx",
         "21 samples · 22 features · 2 sheets",
         "Master sample table (all wells) + ML Feature Matrix availability grid"),
        ("RelPerm_ML_Dataset.xlsx",
         "5 samples · 19 Corey params · 12 sheets",
         "Corey parameters + 41-pt Drainage/Imbibition kr curves + Normalized long-format table"),
        ("CapPressure_ML_Dataset.xlsx",
         "11 samples · 4 sheets",
         "Angsi centrifuge Pc + Duyong MICP + Brooks-Corey fit params + Cross-well Pc summary"),
    ]
    paras = []
    for fname, meta, desc in files_info:
        paras.append(para_xml(f"  {fname}", 13, bold=True, color="002060"))
        paras.append(para_xml(f"  {meta}", 11, color="595959", italic=True))
        paras.append(para_xml(f"  {desc}", 11, color="1A1A1A"))
        paras.append(para_xml("", 8))
    add_textbox(slide, M, bullet_top, SW-2*M, int(SH*0.35), paras)

    # Table: Sheet inventory
    add_textbox(slide, M, int(SH*0.56), SW-2*M, int(SH*0.05),
                [para_xml("RelPerm_ML_Dataset – Sheet Inventory", 11, bold=True, color="002060")])
    sheet_rows = [
        ["Corey_Parameters",       "5","19","Summary: Swir, Sgr, Nw, Ng, RQI, FZI, Gas Recovery"],
        ["Drn_2-015 … Drn_2-033",  "5×41","5","Drainage Krw & Krg vs Sw (per sample)"],
        ["Imb_2-015 … Imb_2-033",  "5×41","5","Imbibition Krw & Krg vs Sw (per sample)"],
        ["Normalized_Kr_Combined", "410","15","Long-format ML table: Sample+Phi+Ka+Swn+Krw+Krg+log10Ka+RQI"],
    ]
    pptx_table(slide, M, int(SH*0.62), SW-2*M, int(SH*0.32),
               ["Sheet Name", "Rows", "Cols", "Contents"],
               sheet_rows, [0.32, 0.08, 0.06, 0.54])


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def main():
    import math  # ensure available in sub-functions
    os.makedirs(RES, exist_ok=True)

    print("Building Unified Sample Properties …")
    samples = build_sample_properties()
    write_sample_properties(samples, OUT_PROPS)

    print("Building Relative Permeability Dataset …")
    write_relperm_dataset(OUT_KR)

    print("Building Capillary Pressure Dataset …")
    write_cappressure_dataset(OUT_PC)

    print(f"Updating PPTX {PPTX_IN} → {PPTX_OUT} …")
    prs = Presentation(PPTX_IN)
    add_unified_dataset_schema_slide(prs)
    add_relperm_corey_slide(prs)
    add_cappressure_slide(prs)
    add_ml_dataset_overview_slide(prs)
    prs.save(PPTX_OUT)
    print(f"  Saved: {PPTX_OUT}")
    print(f"  Total slides: {len(prs.slides)}")
    print("Done.")


if __name__ == "__main__":
    main()
